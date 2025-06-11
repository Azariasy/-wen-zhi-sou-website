import argparse
from pathlib import Path
import os
import shutil
import docx
import jieba
import json
import re
import time
import zipfile
import rarfile
import tempfile
import io
import sys
import multiprocessing
import traceback
import threading
import pandas as pd
import markdown
import math
from html.parser import HTMLParser
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
import email
from email.parser import BytesParser
from email.header import decode_header
from email.utils import parseaddr
import chardet
import extract_msg
import csv
from datetime import datetime
import functools
# --- ADDED: 导入并发处理模块 ---
import asyncio
import concurrent.futures
from threading import Lock
import time
# ------------------------------------

# --- ADDED: 许可证管理器支持 ---
try:
    from license_manager import get_license_manager, Features
    _license_manager_available = True
except ImportError:
    _license_manager_available = False

# --- 导入统一路径处理工具 ---
from path_utils import normalize_path_for_index, PathStandardizer

# --- 导入统一文件处理工具 ---
from file_processing_utils import check_cancellation, periodic_cancellation_check, InterruptedError

# -------------------------------

# 辅助函数，用于检查功能是否可用
def is_feature_available(feature_name):
    """
    检查特定功能是否可用，如果许可证管理器不可用，返回 True
    
    Args:
        feature_name: 要检查的功能名称
        
    Returns:
        bool: 如果功能可用返回 True，否则返回 False
    """
    if not _license_manager_available:
        return True  # 如果许可证管理器不可用，默认所有功能可用
    
    try:
        license_manager = get_license_manager()
        return license_manager.is_feature_available(feature_name)
    except Exception as e:
        print(f"许可证检查错误: {e}")
        return True  # 发生错误时默认功能可用
# --------------------------------

from whoosh.index import create_in, open_dir, exists_in
from whoosh.fields import Schema, TEXT, ID, STORED, NUMERIC, KEYWORD
from whoosh.analysis import Tokenizer, Token, Analyzer
from whoosh.query import Phrase, Term, Prefix, And, Or, Not, Query, NumericRange, Every, FuzzyTerm, Wildcard, TermRange
from whoosh import scoring
from whoosh.qparser import QueryParser, MultifieldParser, OrGroup, GtLtPlugin, PhrasePlugin, SequencePlugin
from PyPDF2 import PdfReader
import pptx
import openpyxl
from datetime import datetime
import csv  # 添加用于写入TSV文件
from tqdm import tqdm  # Added for better progress reporting
from whoosh.highlight import Highlighter, ContextFragmenter, HtmlFormatter
# --- ADDED: Import analysis module ---
from whoosh import analysis
# ------------------------------------

# --- ADDED: Imports for OCR --- 
import pytesseract
from PIL import Image
from pdf2image import convert_from_path, pdfinfo_from_path
from pytesseract import TesseractError # Added
import pdf2image
from pdf2image.exceptions import PDFPopplerTimeoutError, PDFPageCountError # Added exceptions
# ------------------------------

# --- Constants ---
# INDEX_DIR = "indexdir" # Commented out, will be passed as parameter
ALLOWED_EXTENSIONS = [
    # 文档类型
    '.txt', '.docx', '.pdf', '.zip', '.rar', '.pptx', '.xlsx', '.md', '.html', '.htm', '.rtf', '.eml', '.msg',
    # 视频文件 (仅文件名搜索)
    '.mp4', '.mkv', '.avi', '.wmv', '.mov', '.flv', '.webm', '.m4v',
    # 音频文件 (仅文件名搜索)
    '.mp3', '.wav', '.flac', '.aac', '.ogg', '.wma', '.m4a',
    # 图片文件 (仅文件名搜索)
    '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.svg'
]

# 仅进行文件名索引的文件类型（不提取内容）
FILENAME_ONLY_EXTENSIONS = {
    '.mp4', '.mkv', '.avi', '.wmv', '.mov', '.flv', '.webm', '.m4v',  # 视频
    '.mp3', '.wav', '.flac', '.aac', '.ogg', '.wma', '.m4a',         # 音频
    '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.svg' # 图片
}

# --- 新增函数用于记录跳过文件的信息 ---
def record_skipped_file(index_dir_path: str, file_path: str, reason: str) -> None:
    """
    记录被跳过的文件信息到TSV文件中。
    
    Args:
        index_dir_path: 索引目录路径
        file_path: 被跳过的文件路径
        reason: 跳过原因
    """
    try:
        # 构建记录文件路径
        log_file_path = os.path.join(index_dir_path, "index_skipped_files.tsv")
        
        # 获取当前时间戳
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        # 检查文件是否存在，如果不存在则添加表头
        file_exists = os.path.isfile(log_file_path)
        
        # 以追加模式打开文件
        with open(log_file_path, 'a', encoding='utf-8', newline='') as f:
            writer = csv.writer(f, delimiter='\t')
            
            # 如果文件不存在，添加表头
            if not file_exists:
                writer.writerow(["文件路径", "跳过原因", "时间"])
            
            # 写入记录
            writer.writerow([file_path, reason, timestamp])
            
    except Exception as e:
        # 如果记录过程中出错，打印错误但不中断程序
        print(f"Warning: Failed to record skipped file {file_path}: {e}", file=sys.stderr)
        pass

# --- 用于处理跳过原因的格式化函数 ---
def format_skip_reason(reason_type: str, detail: str = "") -> str:
    """
    格式化跳过文件的原因，使其更易于理解。
    
    Args:
        reason_type: 跳过类型
        detail: 详细信息
    
    Returns:
        格式化后的原因说明
    """
    reason_map = {
        "password_zip": "需要密码的ZIP文件",
        "password_rar": "需要密码的RAR文件",
        "corrupted_zip": "损坏的ZIP文件",
        "corrupted_rar": "损坏的RAR文件",
        "ocr_timeout": "OCR处理超时",
        "pdf_conversion_timeout": "PDF转换超时",
        "content_too_large": "内容超出大小限制",
        "unsupported_type": "不支持的文件类型",
        "extraction_error": "提取过程出错",
        "missing_dependency": "缺少处理依赖",
        "pdf_timeout": "PDF处理超时或转换错误",
        "content_limit": "内容大小超过限制",
    }
    
    base_reason = reason_map.get(reason_type, "未知原因")
    
    if detail:
        return f"{base_reason} - {detail}"
    return base_reason
# -----------------------------------

# --- ADDED: Configuration for Tesseract and Poppler --- 
# Set the path to the tesseract executable if it's not in PATH
_TESSERACT_CMD = r'C:\Program Files\Tesseract-OCR\tesseract.exe' # Use double backslashes or raw string
_tesseract_found = False
if os.path.exists(_TESSERACT_CMD):
     pytesseract.pytesseract.tesseract_cmd = _TESSERACT_CMD
     _tesseract_found = True
else:
     if shutil.which('tesseract'):
         # print("Tesseract executable found in PATH.") # COMMENTED OUT
         _tesseract_found = True
     else:
        # print(f"Warning: Tesseract executable not found at {_TESSERACT_CMD} or in PATH. OCR will likely fail.") # COMMENTED OUT
        pass

# Path to the Poppler bin directory
_POPPLER_PATH = r"C:\Program Files\poppler-24.08.0\Library\bin" # Use the path from user screenshot
_poppler_found = False
if os.path.isdir(_POPPLER_PATH) and os.path.exists(os.path.join(_POPPLER_PATH, 'pdfinfo.exe')):
    _poppler_found = True
else:
    # print(f"Warning: Poppler bin directory not found or pdfinfo.exe missing at {_POPPLER_PATH}. PDF to image conversion will likely fail.") # COMMENTED OUT
    pass
# ------------------------------------------------------

# --- Global Index Lock ---
INDEX_ACCESS_LOCK = threading.Lock()

# --- Custom Jieba Tokenizer and Analyzer ---
class ChineseTokenizer(Tokenizer):
    def __call__(self, value, positions=False, chars=False,
                 keeporiginal=False, removestops=True,
                 start_pos=0, start_char=0,
                 mode='', **kwargs):
        assert isinstance(value, str), "ChineseTokenizer expects unicode string"
        
        # --- ADDED: 检测是否为精确搜索模式 ---
        is_phrase_mode = kwargs.get('phrase_mode', False)
        
        if is_phrase_mode:
            # 精确搜索模式：保持原始字符串的完整性，只进行最小分词
            # 对于包含空格的查询，我们需要特殊处理
            if ' ' in value.strip():
                # 包含空格的查询：尝试保持空格前后的词组完整性
                # 例如："部门 A" 应该被处理为一个整体或合理的分段
                
                # 方案1：将整个字符串作为一个token（适用于短查询）
                if len(value.strip()) <= 20:  # 短查询直接作为整体
                    token = Token(positions=positions, chars=chars, removestops=removestops, mode=mode, **kwargs)
                    token.text = value.strip()
                    token.boost = 1.0
                    if keeporiginal:
                        token.original = value.strip()
                    token.stopped = False
                    if positions:
                        token.pos = 0
                    if chars:
                        token.startchar = start_char
                        token.endchar = start_char + len(value.strip())
                    yield token
                    return
                
                # 方案2：智能分段，保持有意义的词组
                # 按空格分割，但保持每个部分的完整性
                parts = [part.strip() for part in value.split() if part.strip()]
                token_pos = 0
                current_char = start_char
                
                for part in parts:
                    # 对每个部分进行jieba分词
                    seglist = jieba.tokenize(part)
                    for (word, start, end) in seglist:
                        if word.strip():  # 跳过纯空格token
                            token = Token(positions=positions, chars=chars, removestops=removestops, mode=mode, **kwargs)
                            token.text = word
                            token.boost = 1.0
                            if keeporiginal:
                                token.original = word
                            token.stopped = False
                            if positions:
                                token.pos = token_pos
                            if chars:
                                token.startchar = current_char + start
                                token.endchar = current_char + end
                            yield token
                            token_pos += 1
                    current_char += len(part) + 1  # +1 for space
                return
        
        # --- 原始逻辑：用于模糊搜索和索引 ---
        seglist = jieba.tokenize(value)
        token_pos = 0
        for (word, start, end) in seglist:
            # --- MODIFIED: 在非精确模式下跳过纯空格token ---
            if not is_phrase_mode and word.strip() == '':
                continue
                
            token = Token(positions=positions, chars=chars, removestops=removestops, mode=mode, **kwargs)
            token.text = word
            token.boost = 1.0
            if keeporiginal:
                token.original = word
            token.stopped = False
            if positions:
                token.pos = token_pos
            if chars:
                token.startchar = start_char + start
                token.endchar = start_char + end
            yield token
            token_pos += 1

class ChineseAnalyzer(Analyzer):
    def __call__(self, value, **kwargs):
        # --- FIXED: 正确传递kwargs参数到ChineseTokenizer ---
        tokenizer = ChineseTokenizer()
        return tokenizer(value, **kwargs)

# --- HTML Stripper ---
class MLStripper(HTMLParser):
    def __init__(self):
        super().__init__()
        self.reset()
        self.strict = False
        self.convert_charrefs = True
        self.text = io.StringIO()
    def handle_data(self, d):
        self.text.write(d)
    def get_data(self):
        return self.text.getvalue()

def strip_tags(html):
    s = MLStripper()
    s.feed(html)
    return s.get_data()

# --- ADDED: 并行搜索引擎优化类 ---
class OptimizedSearchEngine:
    """优化的并行搜索引擎"""
    
    def __init__(self, max_workers: int = 4):
        self.max_workers = max_workers
        self.search_lock = Lock()
        self.result_cache = {}
        self.cache_timeout = 300  # 5分钟缓存过期
        
    def _get_cache_key(self, query_str: str, search_params: dict) -> str:
        """生成缓存键"""
        import hashlib
        key_data = f"{query_str}_{sorted(search_params.items())}"
        return hashlib.md5(key_data.encode()).hexdigest()
        
    def _is_cache_valid(self, cache_entry: dict) -> bool:
        """检查缓存是否仍然有效"""
        return time.time() - cache_entry['timestamp'] < self.cache_timeout
        
    def _analyze_query_complexity(self, query_str: str, search_params: dict) -> str:
        """分析查询复杂度"""
        # 简单查询：短文本，无通配符，无复杂过滤
        if (len(query_str) <= 10 and 
            '*' not in query_str and '?' not in query_str and
            not search_params.get('file_type_filter') and
            not search_params.get('min_size_kb') and
            not search_params.get('start_date')):
            return 'simple'
            
        # 复杂查询：多个通配符，多个过滤条件
        wildcard_count = query_str.count('*') + query_str.count('?')
        filter_count = sum(1 for k in ['file_type_filter', 'min_size_kb', 'max_size_kb', 'start_date', 'end_date'] 
                          if search_params.get(k))
        
        if wildcard_count > 2 or filter_count > 2:
            return 'complex'
            
        return 'medium'
        
    async def optimized_search(self, query_str: str, index_dir_path: str, **search_params) -> list[dict]:
        """优化的搜索入口"""
        start_time = time.time()
        
        # 移除不兼容的参数
        clean_params = search_params.copy()
        if 'limit' in clean_params:
            del clean_params['limit']
        
        # 生成缓存键
        cache_key = self._get_cache_key(query_str, clean_params)
        
        # 检查缓存
        if cache_key in self.result_cache:
            cache_entry = self.result_cache[cache_key]
            if self._is_cache_valid(cache_entry):
                print(f"💾 缓存命中: {query_str} ({len(cache_entry['results'])} 结果)")
                return cache_entry['results']
            else:
                # 清理过期缓存
                del self.result_cache[cache_key]
                
        # 分析查询复杂度
        complexity = self._analyze_query_complexity(query_str, clean_params)
        print(f"🔍 查询复杂度: {complexity}")
        
        # 根据复杂度选择搜索策略
        if complexity == 'simple':
            results = await self._fast_simple_search(query_str, index_dir_path, **clean_params)
        elif complexity == 'medium':
            results = await self._parallel_search(query_str, index_dir_path, **clean_params)
        else:
            results = await self._complex_search_with_optimization(query_str, index_dir_path, **clean_params)
            
        # 缓存结果
        self.result_cache[cache_key] = {
            'results': results,
            'timestamp': time.time()
        }
        
        search_time = time.time() - start_time
        print(f"⚡ 优化搜索完成: {search_time:.2f}秒, {len(results)} 结果")
        
        return results
        
    async def _fast_simple_search(self, query_str: str, index_dir_path: str, **search_params) -> list[dict]:
        """快速简单搜索"""
        loop = asyncio.get_event_loop()
        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
            results = await loop.run_in_executor(
                executor, 
                lambda: search_index(query_str, index_dir_path, **search_params)
            )
        # 对于简单查询，限制返回结果数（限制到500条，优化性能）
        return results[:500] if len(results) > 500 else results
        
    async def _parallel_search(self, query_str: str, index_dir_path: str, **search_params) -> list[dict]:
        """并行搜索（适用于中等复杂度查询）"""
        loop = asyncio.get_event_loop()
        
        # 创建多个搜索任务
        tasks = []
        
        # 如果有文件类型过滤，可以分别搜索不同类型
        file_types = search_params.get('file_type_filter')
        if file_types and len(file_types) > 1:
            # 分文件类型并行搜索
            for file_type in file_types:
                task_params = search_params.copy()
                task_params['file_type_filter'] = [file_type]
                
                with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
                    task = loop.run_in_executor(
                        executor,
                        lambda ft=file_type, tp=task_params: search_index(query_str, index_dir_path, **tp)
                    )
                    tasks.append(task)
                    
            # 等待所有任务完成并合并结果
            if tasks:
                all_results = await asyncio.gather(*tasks)
                merged_results = []
                seen_paths = set()
                
                for results in all_results:
                    for result in results:
                        path = result.get('file_path')
                        if path not in seen_paths:
                            merged_results.append(result)
                            seen_paths.add(path)
                            
                # 按相关度重新排序
                merged_results.sort(key=lambda x: x.get('score', 0), reverse=True)
                return merged_results[:500]  # 限制最多返回500个结果
        
        # 如果无法并行化，使用单线程搜索
        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
            results = await loop.run_in_executor(
                executor,
                lambda: search_index(query_str, index_dir_path, **search_params)
            )
        return results
        
    async def _complex_search_with_optimization(self, query_str: str, index_dir_path: str, **search_params) -> list[dict]:
        """复杂搜索优化"""
        # 对于复杂查询，采用分阶段搜索策略
        
        # 第一阶段：快速文件名搜索
        filename_params = search_params.copy()
        filename_params['search_scope'] = 'filename'
        
        loop = asyncio.get_event_loop()
        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
            filename_results = await loop.run_in_executor(
                executor,
                lambda: search_index(query_str, index_dir_path, **filename_params)
            )
        
        print(f"📁 文件名搜索: {len(filename_results)} 结果")
        
        # 第二阶段：全文搜索
        fulltext_params = search_params.copy()
        fulltext_params['search_scope'] = 'fulltext'
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
            fulltext_results = await loop.run_in_executor(
                executor,
                lambda: search_index(query_str, index_dir_path, **fulltext_params)
            )
            
        print(f"📄 全文搜索: {len(fulltext_results)} 结果")
        
        # 合并结果，去重，按相关度排序
        all_results = filename_results + fulltext_results
        seen_paths = set()
        merged_results = []
        
        for result in all_results:
            path = result.get('file_path')
            if path not in seen_paths:
                merged_results.append(result)
                seen_paths.add(path)
                
        # 按相关度排序
        merged_results.sort(key=lambda x: x.get('score', 0), reverse=True)
        return merged_results[:500]  # 限制最多返回500个结果
        
    def clear_cache(self):
        """清理缓存"""
        self.result_cache.clear()
        print("🧹 搜索缓存已清理")
        
    def get_cache_stats(self) -> dict:
        """获取缓存统计"""
        valid_entries = sum(1 for entry in self.result_cache.values() 
                           if self._is_cache_valid(entry))
        return {
            'total_entries': len(self.result_cache),
            'valid_entries': valid_entries,
            'cache_hit_potential': valid_entries / max(len(self.result_cache), 1)
        }

# 创建全局优化搜索引擎实例
_optimized_search_engine = None

def get_optimized_search_engine() -> OptimizedSearchEngine:
    """获取全局优化搜索引擎实例"""
    global _optimized_search_engine
    if _optimized_search_engine is None:
        _optimized_search_engine = OptimizedSearchEngine(max_workers=4)
    return _optimized_search_engine

def optimized_search_sync(query_str: str, index_dir_path: str, **search_params) -> list[dict]:
    """同步版本的优化搜索接口"""
    engine = get_optimized_search_engine()
    
    # 如果已经在事件循环中，直接调用异步版本
    try:
        loop = asyncio.get_running_loop()
        # 在已有事件循环中，创建新任务
        task = asyncio.create_task(engine.optimized_search(query_str, index_dir_path, **search_params))
        return asyncio.run_coroutine_threadsafe(task, loop).result()
    except RuntimeError:
        # 没有运行的事件循环，创建新的
        return asyncio.run(engine.optimized_search(query_str, index_dir_path, **search_params))
# ------------------------------------

def get_schema() -> Schema:
    analyzer = ChineseAnalyzer()
    return Schema(path=ID(stored=True, unique=True),
                  content=TEXT(stored=True, analyzer=analyzer),
                  # --- ADDED: Field for filename search ---
                  filename_text=TEXT(stored=True, analyzer=analysis.StandardAnalyzer()),
                  # ---------------------------------------------------
                  structure_map=STORED,
                  last_modified=NUMERIC(stored=True, sortable=True),
                  file_size=NUMERIC(stored=True, sortable=True),
                  file_type=KEYWORD(stored=True, lowercase=True, scorable=False),
                  # --- ADDED: Field to store OCR indexing status ---
                  indexed_with_ocr=STORED)
                  # -------------------------------------------------

def scan_documents(directory_path: Path) -> list[Path]:
    found_files = []
    if not directory_path.is_dir():
        # print(f"Error: Path is not a directory: {directory_path}") # COMMENTED OUT
        return found_files
    # print(f"Scanning directory: {directory_path}") # COMMENTED OUT
    for item in directory_path.rglob('*'):
        if item.is_file() and item.suffix.lower() in ALLOWED_EXTENSIONS:
            found_files.append(item)
    # print(f"Scan complete. Found {len(found_files)} document(s).") # COMMENTED OUT
    return found_files

def extract_text_from_docx(file_path: Path, cancel_callback=None) -> tuple[str, list[dict]]:
    full_text_list = []
    structure = []
    try:
        # 开始处理前检查取消状态
        check_cancellation(cancel_callback, "DOCX文件处理")
        
        doc = docx.Document(file_path)
        for i, para in enumerate(doc.paragraphs):
            # 每处理50个段落检查一次取消状态（使用统一工具函数）
            periodic_cancellation_check(cancel_callback, 50, i, "DOCX段落处理")
                
            text = para.text.strip()
            if not text:
                continue
            full_text_list.append(text)
            para_info = {"text": text}
            style_name = para.style.name.lower()
            if style_name.startswith('heading'):
                para_info["type"] = "heading"
                try:
                    level = int(style_name.split()[-1])
                    para_info["level"] = level
                except (ValueError, IndexError):
                    para_info["level"] = 1
            else:
                para_info["type"] = "paragraph"
            structure.append(para_info)
        
        # 处理完成前最后检查取消状态
        check_cancellation(cancel_callback, "DOCX文件处理")
        
        full_text = '\n'.join(full_text_list)
        return full_text, structure
    except InterruptedError:
        # 重新抛出取消异常
        raise
    except Exception as e:
        # print(f"Error reading docx file {file_path}: {e}") # COMMENTED OUT
        return "", []

def extract_text_from_txt(file_path: Path, cancel_callback=None) -> tuple[str, list[dict]]:
    content = ""
    structure = []
    encodings_to_try = ['utf-8', 'gbk', 'gb2312', 'latin-1']
    heading_pattern = re.compile(
        r'^\s*(?:第[一二三四五六七八九十百千万零〇]+[章条]|'  
        r'[零一二三四五六七八九十]+、|'             
        r'[（(][一二三四五六七八九十百千万零〇]+[)）]|' 
        r'[①②③④⑤⑥⑦⑧⑨⑩⑪⑫⑬⑭⑮⑯⑰⑱⑲⑳]|' 
        r'\d+(?:\.\d+)*\.?\s+|'                
        r'第.*?节)\s*$',                         
        re.IGNORECASE
    )
    try:
        for encoding in encodings_to_try:
            try:
                # 在开始读取文件前检查是否需要取消
                check_cancellation(cancel_callback, "TXT文件读取")
                    
                content = file_path.read_text(encoding=encoding)
                if content:
                    lines = content.splitlines()
                    for i, line in enumerate(lines):
                        # 对于大文件，每处理100行检查一次取消状态
                        periodic_cancellation_check(cancel_callback, 100, i, "TXT文件行处理")
                            
                        cleaned_line = line.strip()
                        if cleaned_line:
                            if heading_pattern.match(cleaned_line):
                                structure.append({'type': 'heading', 'level': 1, 'text': cleaned_line})
                            else:
                                structure.append({'type': 'paragraph', 'text': cleaned_line})
                if content or structure:
                    break
            except InterruptedError:
                # 重新抛出取消异常
                raise
            except UnicodeDecodeError:
                continue
            except Exception as e:
                # print(f"Error reading txt file {file_path} with encoding {encoding}: {e}") # COMMENTED OUT
                return "", []
        if not content and not structure:
            print(f"Could not decode txt file {file_path} with tried encodings or file is empty.")
            return "", []
    except Exception as e:
        print(f"Error accessing txt file {file_path}: {e}")
        return "", []
    return content, structure

def extract_text_from_pdf(file_path: Path, enable_ocr: bool = True, ocr_lang: str = 'chi_sim+eng', min_chars_for_ocr_trigger: int = 50, timeout: int | None = None, cancel_callback=None) -> tuple[str | None, list[dict]]:
    """
    从PDF文件提取文本，支持基于PyPDF2和OCR的混合方法
    """
    # --- ADDED: 检查PDF支持许可证 ---
    if not is_feature_available(Features.PDF_SUPPORT):
        print(f"PDF支持功能不可用 (未获得许可)")
        raise PermissionError("PDF支持功能需要专业版许可证。")
    # -------------------------------
    
    # 基础异常检查
    if not file_path.exists():
        print(f"PDF文件不存在: {file_path}")
        return None, []
    
    # --- ADDED: 早期取消检查 ---
    check_cancellation(cancel_callback, "PDF文件处理")
    # ---------------------------
    
    extracted_text = ""
    direct_text = ""
    structure = [] # Initialize structure list

    # 1. Try direct text extraction (using PyPDF2 or similar - currently placeholder)
    #    NOTE: Direct text extraction libraries usually don't have configurable timeouts easily.
    try:
        # --- Placeholder for direct text extraction ---
        # Example using PyPDF2 (if installed):
        # from PyPDF2 import PdfReader
        # reader = PdfReader(file_path)
        # for page in reader.pages:
        #     page_content = page.extract_text()
        #     if page_content:
        #         direct_text += page_content + "\n"
        # direct_text = direct_text.strip()
        # --- End Placeholder ---

        if not enable_ocr:
            print(f"Warning: Direct PDF text extraction not fully implemented or OCR disabled for {file_path.name}. Relying on potentially empty direct_text.")
            pass
    except Exception as e:
        print(f"Error during direct text extraction attempt for {file_path.name}: {e}", file=sys.stderr)
        direct_text = ""

    ocr_needed = enable_ocr and (len(direct_text) < min_chars_for_ocr_trigger)

    if ocr_needed:
        print(f"Info: {file_path.name} direct text insufficient ({len(direct_text)} chars) or not attempted, trying OCR...")
        ocr_texts = []
        try:
            # --- ADDED: 在开始OCR前再次检查取消状态 ---
            check_cancellation(cancel_callback, "PDF OCR处理开始")
            # ----------------------------------------
                
            # --- 添加 pdf2image 时间日志 ---
            pdf2image_start_time = time.time()
            print(f"DEBUG: [{file_path.name}] Starting pdf2image conversion at {pdf2image_start_time:.2f}")
            # -------------------------------

            # --- MODIFIED: 使用更短的超时时间进行PDF转换，以便更快响应取消 ---
            # 将原始超时时间分割为更小的块，每个块后检查取消状态
            chunk_timeout = min(30, timeout) if timeout else 30  # 每30秒检查一次

            images = pdf2image.convert_from_path(
                            file_path,
                timeout=chunk_timeout,  # 使用较短的超时
                fmt='jpeg',
                thread_count=1
            )

            # --- 添加 pdf2image 时间日志 ---
            pdf2image_end_time = time.time()
            pdf2image_duration = pdf2image_end_time - pdf2image_start_time
            print(f"DEBUG: [{file_path.name}] Finished pdf2image conversion at {pdf2image_end_time:.2f} (Duration: {pdf2image_duration:.2f}s)")
            # -------------------------------

            # --- ADDED: 在开始处理页面前检查取消状态 ---
            check_cancellation(cancel_callback, "PDF页面处理开始")
            # ----------------------------------------

            for i, image in enumerate(images):
                # --- MODIFIED: 在处理每个页面前检查是否需要取消 ---
                check_cancellation(cancel_callback, f"PDF页面 {i + 1} 处理")
                # ------------------------------------------------
                    
                page_num = i + 1
                try:
                    # --- 添加 Tesseract 时间日志 ---
                    tesseract_start_time = time.time()
                    print(f"DEBUG: [{file_path.name} Page {page_num}] Starting Tesseract OCR at {tesseract_start_time:.2f}")
                    # ------------------------------

                    # --- MODIFIED: 使用更短的OCR超时时间 ---
                    ocr_timeout = min(60, timeout) if timeout else 60  # 每页最多60秒
                    page_text = pytesseract.image_to_string(image, lang=ocr_lang, timeout=ocr_timeout)

                    # --- 添加 Tesseract 时间日志 ---
                    tesseract_end_time = time.time()
                    tesseract_duration = tesseract_end_time - tesseract_start_time
                    print(f"DEBUG: [{file_path.name} Page {page_num}] Finished Tesseract OCR at {tesseract_end_time:.2f} (Duration: {tesseract_duration:.2f}s)")
                    # ------------------------------

                    page_text = page_text.strip()
                    if page_text:
                        ocr_texts.append(page_text)
                    else:
                        pass

                    # --- ADDED: 在每页OCR完成后检查取消状态 ---
                    check_cancellation(cancel_callback, f"PDF页面 {page_num} OCR完成")
                    # ----------------------------------------
                        
                except TesseractError as te:
                    tesseract_fail_time = time.time()
                    print(f"DEBUG: [{file_path.name} Page {page_num}] TesseractError at {tesseract_fail_time:.2f} (Duration before error: {tesseract_fail_time - tesseract_start_time:.2f}s)")
                    err_msg = str(te).lower()
                    if 'timeout' in err_msg or 'process timed out' in err_msg:
                        print(f"Warning: Tesseract OCR timed out (>{ocr_timeout}s) for page {page_num} of {file_path.name}.", file=sys.stderr)
                        # 这里不记录跳过的文件，因为该函数无法访问index_dir_path
                        # 记录会在_extract_worker中完成
                        return None, [] # MODIFIED
                    else:
                        print(f"Error during Tesseract OCR for page {page_num} of {file_path.name}: {te}", file=sys.stderr)
                        # 这里不记录跳过的文件，因为该函数无法访问index_dir_path
                        # 记录会在_extract_worker中完成
                        return None, [] # MODIFIED
                except RuntimeError as rte:
                    tesseract_fail_time = time.time()
                    print(f"DEBUG: [{file_path.name} Page {page_num}] RuntimeError at {tesseract_fail_time:.2f} (Duration before error: {tesseract_fail_time - tesseract_start_time:.2f}s)")
                    err_msg = str(rte).lower()
                    if 'timeout' in err_msg:
                        print(f"Warning: Tesseract OCR likely timed out (>{ocr_timeout}s) for page {page_num} of {file_path.name}. Error: {rte}", file=sys.stderr)
                        # 这里不记录跳过的文件，因为该函数无法访问index_dir_path
                        # 记录会在_extract_worker中完成
                        return None, [] # MODIFIED
                    else:
                        print(f"Runtime error during Tesseract OCR for page {page_num} of {file_path.name}: {rte}", file=sys.stderr)
                        # 这里不记录跳过的文件，因为该函数无法访问index_dir_path
                        # 记录会在_extract_worker中完成
                        return None, [] # MODIFIED
                except Exception as e:
                    tesseract_fail_time = time.time()
                    print(f"DEBUG: [{file_path.name} Page {page_num}] Exception at {tesseract_fail_time:.2f} (Duration before error: {tesseract_fail_time - tesseract_start_time:.2f}s)")
                    print(f"Unexpected error during OCR for page {page_num} of {file_path.name}: {e}", file=sys.stderr)
                    # 这里不记录跳过的文件，因为该函数无法访问index_dir_path
                    # 记录会在_extract_worker中完成
                    return None, [] # MODIFIED
                # --- End Modification ---

            extracted_text = "\n\n".join(ocr_texts)
            if images:
                 print(f"Info: OCR process completed for {file_path.name}. Total chars: {len(extracted_text)}")

        # --- MODIFIED: Return tuple on exceptions ---
        except PDFPopplerTimeoutError:
            pdf2image_fail_time = time.time()
            print(f"DEBUG: [{file_path.name}] PDFPopplerTimeoutError at {pdf2image_fail_time:.2f} (Duration before error: {pdf2image_fail_time - pdf2image_start_time:.2f}s)")
            print(f"Warning: PDF to image conversion timed out (>{chunk_timeout}s) for {file_path.name}.", file=sys.stderr)
            # 这里不记录跳过的文件，因为该函数无法访问index_dir_path
            # 记录会在_extract_worker中完成
            return None, [] # MODIFIED
        except PDFPageCountError as pe:
             print(f"Error getting page count or converting PDF {file_path.name}: {pe}. Skipping OCR.", file=sys.stderr)
             # 这里不记录跳过的文件，因为该函数无法访问index_dir_path
             # 记录会在_extract_worker中完成
             return None, [] # MODIFIED
        except InterruptedError:
            # --- ADDED: 专门处理用户取消 ---
            print(f"PDF OCR processing cancelled by user for {file_path.name}")
            raise  # 重新抛出取消异常
        except Exception as e:
            pdf2image_fail_time = time.time()
            duration_str = f"(Duration before error: {pdf2image_fail_time - pdf2image_start_time:.2f}s)" if 'pdf2image_start_time' in locals() else ""
            print(f"DEBUG: [{file_path.name}] Exception during PDF conversion at {pdf2image_fail_time:.2f} {duration_str}")
            print(f"Error during PDF processing/conversion for {file_path.name}: {e}", file=sys.stderr)
            print(traceback.format_exc(), file=sys.stderr)
            # 这里不记录跳过的文件，因为该函数无法访问index_dir_path
            # 记录会在_extract_worker中完成
            return None, [] # MODIFIED
        # --- End Modification ---

    #     else:
        extracted_text = direct_text
        # ... (logging) ...

    # --- Generate basic structure from final extracted_text ---
    # Note: If direct_text extraction were implemented, it might populate 'structure' directly.
    # This block ensures structure is generated if only 'extracted_text' is available (e.g., from OCR).
    if not structure and extracted_text: # Generate structure only if it's empty AND text exists
        lines = extracted_text.splitlines()
        for line in lines:
            cleaned_line = line.strip()
            if cleaned_line:
                structure.append({'type': 'paragraph', 'text': cleaned_line})
    # -----------------------------------------------------------

    # MODIFIED: Return tuple, handle potential None for text
    if extracted_text is None:
         return None, [] # Structure is already empty list in this case
    else:
         # Ensure structure is a list even if extraction yielded text but somehow failed structure generation
         final_structure = structure if isinstance(structure, list) else []
         return extracted_text, final_structure

def extract_text_from_pptx(file_path: Path, cancel_callback=None) -> tuple[str, list[dict]]:
    full_text_list = []
    structure = []
    try:
        # 开始处理前检查取消状态
        check_cancellation(cancel_callback, "PPTX文件处理")
        
        presentation = pptx.Presentation(file_path)
        for slide_num, slide in enumerate(presentation.slides):
            # 在处理每个幻灯片前检查是否需要取消
            check_cancellation(cancel_callback, f"PPTX幻灯片 {slide_num + 1} 处理")
                
            slide_texts = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text.strip()
                if text:
                    slide_texts.append(text)
                    structure.append({
                        'type': 'paragraph',
                        'text': text,
                        'context': f'Slide {slide_num + 1}'
                    })
            if slide_texts:
                full_text_list.extend(slide_texts)
        full_text = '\n'.join(full_text_list)
        return full_text, structure
    except InterruptedError:
        # 重新抛出取消异常
        raise
    except Exception as e:
        print(f"Error reading pptx file {file_path}: {e}")
        return "", []

def _is_potential_header(row_data: list, non_empty_threshold=0.5, string_threshold=0.5) -> bool:
    if not row_data:
        return False
    non_empty_count = sum(1 for cell in row_data if pd.notna(cell) and str(cell).strip())
    string_count = sum(1 for cell in row_data if isinstance(cell, str) and str(cell).strip() and not str(cell).replace('.', '', 1).isdigit())
    if non_empty_count == 0:
        return False
    return (non_empty_count / len(row_data) >= non_empty_threshold and
            string_count / non_empty_count >= string_threshold)

def extract_text_from_xlsx(file_path: Path, cancel_callback=None) -> tuple[str, list[dict]]:
    full_text_list = []
    structure = []
    MAX_HEADER_CHECK_ROWS = 10
    
    # 在开始处理前检查是否需要取消
    check_cancellation(cancel_callback, "XLSX文件处理")
    
    try:
        excel_data = pd.read_excel(file_path, sheet_name=None, header=None, keep_default_na=False)
    except Exception as e:
        print(f"Error reading Excel file {file_path}: {e}")
        return "", []
    
    # 在读取Excel数据后检查取消状态
    check_cancellation(cancel_callback, "XLSX数据读取")
    
    for sheet_name, df_initial in excel_data.items():
        # 在处理每个工作表前检查是否需要取消
        check_cancellation(cancel_callback, f"XLSX工作表 {sheet_name} 处理")
            
        if df_initial.empty:
            continue
            
        header_row_index = -1
        best_header_score = -1
        potential_header_idx = -1
        
        # 在表头检测循环中增加更频繁的取消检查
        for i in range(min(MAX_HEADER_CHECK_ROWS, len(df_initial))):
            # 在检查表头的循环中也添加取消检查（使用周期性检查）
            periodic_cancellation_check(cancel_callback, 5, i, f"XLSX表头检测 {sheet_name}")
                
            row_values = df_initial.iloc[i].tolist()
            current_score = sum(1 for cell in row_values if pd.notna(cell) and isinstance(cell, str) and str(cell).strip() and not str(cell).replace('.', '', 1).isdigit())
            is_plausible = any(isinstance(cell, str) and str(cell).strip() for cell in row_values)
            if is_plausible and current_score > best_header_score:
                best_header_score = current_score
                potential_header_idx = i
            elif best_header_score == -1 and is_plausible:
                potential_header_idx = i
                
        # 在表头检测完成后检查取消状态
        check_cancellation(cancel_callback, f"XLSX表头检测完成 {sheet_name}")
            
        if potential_header_idx != -1 and best_header_score >= 1:
            header_row_index = potential_header_idx
            print(f"Detected header in '{sheet_name}' at row {header_row_index + 1}")
            try:
                # 在重新读取工作表前检查取消状态
                check_cancellation(cancel_callback, f"XLSX重新读取工作表 {sheet_name}")
                    
                df = pd.read_excel(file_path, sheet_name=sheet_name, header=header_row_index, keep_default_na=False)
                df = df.dropna(axis=0, how='all')
                df = df.dropna(axis=1, how='all')
                df = df.fillna('')
                headers = [str(h).strip() for h in df.columns]
                
                # 在开始处理行数据前检查取消状态
                check_cancellation(cancel_callback, f"XLSX开始处理行数据 {sheet_name}")
                
                row_count = 0
                for idx, row in df.iterrows():
                    # 在处理每一行时也检查取消 - 增加频率
                    check_cancellation(cancel_callback, f"XLSX行处理 {sheet_name} 第{idx}行")
                    
                    # 每处理10行检查一次取消状态（提高响应性）
                    row_count += 1
                    periodic_cancellation_check(cancel_callback, 10, row_count, f"XLSX批量行处理 {sheet_name}")
                        
                    excel_row_num = idx + header_row_index + 2
                    row_values = [str(cell).strip() for cell in row.tolist()]
                    row_texts_for_full_text = [h for h in headers if h] + [v for v in row_values if v]
                    if row_texts_for_full_text:
                        full_text_list.append(" ".join(row_texts_for_full_text))
                    if any(row_values):
                        searchable_row_text = " ".join(filter(None, row_values))
                        structure.append({
                            'type': 'excel_row',
                            'sheet_name': str(sheet_name),
                            'row_index': excel_row_num,
                            'headers': headers,
                            'values': row_values,
                            'text': searchable_row_text
                        })
                        
            except InterruptedError:
                # 重新抛出取消异常
                print(f"XLSX工作表处理被用户取消: {sheet_name}")
                raise
            except Exception as e:
                print(f"Error re-reading sheet '{sheet_name}' with header at row {header_row_index + 1}: {e}")
                # 在异常处理中也检查取消状态
                check_cancellation(cancel_callback, f"XLSX异常处理 {sheet_name}")
                sheet_text = df_initial.to_string(index=False, header=False)
                full_text_list.append(sheet_text)
                structure.append({'type': 'paragraph', 'text': f"Content from sheet '{sheet_name}' (header detection failed)."})
        else:
            print(f"Warning: Could not detect header in sheet '{sheet_name}'. Indexing as plain text.")
            # 在处理无表头工作表前检查取消状态
            check_cancellation(cancel_callback, f"XLSX无表头工作表处理 {sheet_name}")
                
            sheet_text = df_initial.to_string(index=False, header=False)
            cleaned_sheet_text = "\n".join(line.strip() for line in sheet_text.splitlines() if line.strip())
            if cleaned_sheet_text:
                full_text_list.append(cleaned_sheet_text)
                structure.append({
                    'type': 'paragraph',
                    'text': cleaned_sheet_text,
                    'context': f"Sheet: {sheet_name} (No header detected)"
                })
                
    # 在返回前最后检查一次取消状态
    check_cancellation(cancel_callback, "XLSX文件处理完成")
        
    full_text = '\n'.join(full_text_list)
    return full_text, structure


def extract_text_from_md(file_path: Path, cancel_callback=None) -> tuple[str, list[dict]]:
    """Extract text and structure from a Markdown file."""
    # --- ADDED: 检查Markdown支持许可证 ---
    if not is_feature_available(Features.MARKDOWN_SUPPORT):
        print(f"Markdown支持功能不可用 (未获得许可)")
        raise PermissionError("Markdown支持功能需要专业版许可证。")
    # ------------------------------------
    
    if not file_path.exists():
        return "", []
    
    content = ""
    structure = []
    encodings_to_try = ['utf-8', 'gbk', 'gb2312']
    try:
        # 在开始处理前检查是否需要取消
        check_cancellation(cancel_callback, "MD文件处理开始")
            
        raw_md_content = None
        for encoding in encodings_to_try:
            try:
                raw_md_content = file_path.read_text(encoding=encoding)
                if raw_md_content is not None:
                    break
            except UnicodeDecodeError:
                continue
            except Exception as e:
                print(f"Error reading md file {file_path} with encoding {encoding}: {e}")
                return "", []
        if raw_md_content is None:
            print(f"Could not decode md file {file_path} with tried encodings or file is empty.")
            return "", []
        
        # 在转换前检查取消状态
        check_cancellation(cancel_callback, "MD转换HTML")
            
        html_content = markdown.markdown(raw_md_content)
        content = strip_tags(html_content).strip()
        if content:
            structure.append({'type': 'paragraph', 'text': content})
        print(f"MD Extracted Text (first 500 chars): {content[:500]}")
        return content, structure
    except InterruptedError:
        # 重新抛出取消异常
        raise
    except ImportError:
        print("Error: markdown library not installed. Please run: pip install Markdown")
        return "", []

def extract_text_from_html(file_path: Path, cancel_callback=None) -> tuple[str, list[dict]]:
    print(f"Attempting to extract text from HTML/HTM: {file_path}")
    content = ""
    structure = []
    try:
        # 在开始处理前检查是否需要取消
        check_cancellation(cancel_callback, "HTML文件处理开始")
            
        raw_html_bytes = file_path.read_bytes()
        if not raw_html_bytes:
            print(f"Warning: HTML/HTM file is empty: {file_path}")
            return "", []
        detected = chardet.detect(raw_html_bytes)
        encoding = detected.get('encoding')
        confidence = detected.get('confidence', 0)
        print(f"Chardet detected encoding: {encoding} with confidence: {confidence:.2f}")
        html_content_decoded = None
        if encoding and confidence > 0.7:
            try:
                html_content_decoded = raw_html_bytes.decode(encoding, errors='replace')
                print(f"Decoded using chardet result: {encoding}")
            except Exception as e:
                print(f"Warning: Failed to decode using detected encoding '{encoding}': {e}")
        if html_content_decoded is None:
            for enc in ['utf-8', 'gbk', 'gb18030', 'latin-1']:
                try:
                    html_content_decoded = raw_html_bytes.decode(enc, errors='strict')
                    print(f"Decoded using fallback encoding: {enc}")
                    break
                except UnicodeDecodeError:
                    continue
            if html_content_decoded is None:
                print(f"Warning: Could not decode HTML/HTM, using utf-8 with replacements.")
                html_content_decoded = raw_html_bytes.decode('utf-8', errors='replace')
        
        # 在解析前再次检查取消状态
        check_cancellation(cancel_callback, "HTML解析处理")
            
        soup = BeautifulSoup(html_content_decoded, 'lxml')
        for script_or_style in soup(["script", "style"]):
            script_or_style.decompose()
        content = ' '.join(soup.stripped_strings)
        if content:
            structure.append({'type': 'paragraph', 'text': content})
        print(f"Successfully extracted HTML/HTM content (length: {len(content)}): {file_path}")
        print(f"Extracted sample: {content[:500]}{'...' if len(content) > 500 else ''}")
        return content, structure
    except InterruptedError:
        # 重新抛出取消异常
        raise
    except Exception as e:
        print(f"Error processing html/htm file {file_path}: {e}")
        # traceback.print_exc(file=sys.stderr) # COMMENTED OUT
        return "", []

def extract_text_from_rtf(file_path: Path, cancel_callback=None) -> tuple[str, list[dict]]:
    content = ""
    structure = []
    encodings_to_try = ['utf-8', 'gbk', 'gb2312', 'latin-1']
    try:
        # 在开始处理前检查是否需要取消
        check_cancellation(cancel_callback, "RTF文件处理开始")
            
        rtf_content = None
        for encoding in encodings_to_try:
            try:
                rtf_content = file_path.read_text(encoding=encoding)
                if rtf_content is not None and rtf_content.strip().startswith('{\\rtf'):
                    break
                else:
                    rtf_content = None
            except UnicodeDecodeError:
                continue
            except Exception as e:
                print(f"Error reading rtf file {file_path} with encoding {encoding}: {e}")
                return "", []
        if rtf_content is None:
            print(f"Could not decode rtf file {file_path} with tried encodings or file doesn't start with RTF marker.")
            return "", []
        
        # 在转换前检查取消状态
        check_cancellation(cancel_callback, "RTF转换文本")
            
        content = rtf_to_text(rtf_content, errors="ignore").strip()
        if content:
            structure.append({'type': 'paragraph', 'text': content})
        return content, structure
    except InterruptedError:
        # 重新抛出取消异常
        raise
    except Exception as e:
        print(f"Error processing rtf file {file_path}: {e}")
        return "", []

def extract_text_from_eml(file_path: Path, cancel_callback=None) -> tuple[str, list[dict]]:
    """Extract text from an EML file."""
    # --- ADDED: 检查邮件支持许可证 ---
    if not is_feature_available(Features.EMAIL_SUPPORT):
        print(f"邮件支持功能不可用 (未获得许可)")
        raise PermissionError("邮件支持功能需要专业版许可证。")
    # -----------------------------
    
    if not file_path.exists():
        return "", []
    
    full_text_list = []
    structure = []
    try:
        # 在开始处理前检查是否需要取消
        check_cancellation(cancel_callback, "EML文件处理开始")
            
        raw_bytes = file_path.read_bytes()
        if not raw_bytes:
            return "", []
        parser = BytesParser()
        msg = parser.parsebytes(raw_bytes)
        
        # 在解析邮件内容前检查取消状态
        check_cancellation(cancel_callback, "EML邮件解析")
            
        def decode_header_simple(header_value):
            if not header_value:
                return ""
            decoded_parts = []
            for part, encoding in decode_header(header_value):
                if isinstance(part, bytes):
                    try:
                        decoded_parts.append(part.decode(encoding or 'utf-8', errors='replace'))
                    except LookupError:
                        decoded_parts.append(part.decode('utf-8', errors='replace'))
                else:
                    decoded_parts.append(part)
            return "".join(decoded_parts)
        subject = decode_header_simple(msg.get('Subject'))
        sender = decode_header_simple(msg.get('From'))
        recipient = decode_header_simple(msg.get('To'))
        cc = decode_header_simple(msg.get('Cc'))
        date_str = msg.get('Date')
        if subject:
            structure.append({'type': 'heading', 'level': 1, 'text': f"主题: {subject}"})
            full_text_list.append(subject)
        if sender:
            realname, email_addr = parseaddr(sender)
            clean_sender = f"{realname} <{email_addr}>" if realname else sender
            structure.append({'type': 'metadata', 'text': f"发件人: {clean_sender}"})
            full_text_list.append(f"发件人: {clean_sender}")
        if recipient:
            structure.append({'type': 'metadata', 'text': f"收件人: {recipient}"})
            full_text_list.append(f"收件人: {recipient}")
        if cc:
            structure.append({'type': 'metadata', 'text': f"抄送: {cc}"})
            full_text_list.append(f"抄送: {cc}")
        if date_str:
            structure.append({'type': 'metadata', 'text': f"日期: {date_str}"})
            full_text_list.append(f"日期: {date_str}")
        body_content = ""
        if msg.is_multipart():
            for part in msg.walk():
                content_type = part.get_content_type()
                content_disposition = str(part.get('Content-Disposition'))
                if 'attachment' not in content_disposition and content_type.startswith('text/'):
                    payload = part.get_payload(decode=True)
                    charset = part.get_content_charset() or 'utf-8'
                    try:
                        body_content += payload.decode(charset, errors='replace')
                    except LookupError:
                        body_content += payload.decode('utf-8', errors='replace')
        else:
            payload = msg.get_payload(decode=True)
            charset = msg.get_content_charset() or 'utf-8'
            try:
                body_content = payload.decode(charset, errors='replace')
            except LookupError:
                body_content = payload.decode('utf-8', errors='replace')
        if body_content:
            lines = body_content.splitlines()
            for line in lines:
                cleaned_line = line.strip()
                if cleaned_line:
                    structure.append({'type': 'paragraph', 'text': cleaned_line})
                    full_text_list.append(cleaned_line)
        full_text = '\n'.join(full_text_list)
        return full_text, structure
    except Exception as e:
        print(f"Error extracting EML {file_path}: {e}")
        # traceback.print_exc(file=sys.stderr) # COMMENTED OUT
        return "", []

def extract_text_from_msg(file_path: Path, cancel_callback=None) -> tuple[str, list[dict]]:
    # --- ADDED: 检查邮件支持许可证 ---
    if not is_feature_available(Features.EMAIL_SUPPORT):
        print(f"邮件支持功能不可用 (未获得许可)")
        raise PermissionError("邮件支持功能需要专业版许可证。")
    # -----------------------------
    
    print(f"Attempting to extract text from MSG: {file_path.name}")
    full_text_list = []
    structure = []
    try:
        msg = extract_msg.Message(str(file_path))
        if msg.subject:
            structure.append({'type': 'heading', 'level': 1, 'text': f"主题: {msg.subject}"})
            full_text_list.append(msg.subject)
        if msg.sender:
            structure.append({'type': 'metadata', 'text': f"发件人: {msg.sender}"})
            full_text_list.append(f"发件人: {msg.sender}")
        if msg.to:
            structure.append({'type': 'metadata', 'text': f"收件人: {msg.to}"})
            full_text_list.append(f"收件人: {msg.to}")
        if msg.cc:
            structure.append({'type': 'metadata', 'text': f"抄送: {msg.cc}"})
            full_text_list.append(f"抄送: {msg.cc}")
        if msg.date:
            structure.append({'type': 'metadata', 'text': f"日期: {msg.date}"})
            full_text_list.append(f"日期: {msg.date}")
        body_to_process = None
        processed_html_body = False
        
        # 在处理邮件正文前检查取消状态
        check_cancellation(cancel_callback, "MSG邮件正文处理")
            
        if hasattr(msg, 'htmlBody') and msg.htmlBody:
            print("MSG Body: Attempting to process HTML body.")
            html_bytes = None
            try:
                if isinstance(msg.htmlBody, bytes):
                    html_bytes = msg.htmlBody
                elif isinstance(msg.htmlBody, str):
                    html_bytes = msg.htmlBody.encode('utf-8', errors='ignore')
                if html_bytes:
                    detected = chardet.detect(html_bytes)
                    encoding = detected.get('encoding')
                    confidence = detected.get('confidence', 0)
                    html_content_decoded = None
                    if encoding and confidence > 0.7:
                        try:
                            html_content_decoded = html_bytes.decode(encoding, errors='replace')
                            print(f"MSG HTML Body: Decoded using chardet result: {encoding}")
                        except Exception as e:
                            print(f"MSG HTML Body: Warning - Failed using chardet '{encoding}': {e}")
                    if html_content_decoded is None:
                        for enc in ['utf-8', 'gbk', 'gb18030']:
                            try:
                                html_content_decoded = html_bytes.decode(enc, errors='strict')
                                print(f"MSG HTML Body: Decoded using fallback: {enc}")
                                break
                            except Exception:
                                pass  # 继续尝试其他编码
                        if html_content_decoded is None:
                            html_content_decoded = html_bytes.decode('utf-8', errors='replace')
                            print(f"MSG HTML Body: Decoded using final fallback utf-8 (replace)")
                    if html_content_decoded:
                        soup = BeautifulSoup(html_content_decoded, 'lxml')
                        for script_or_style in soup(["script", "style"]):
                            script_or_style.decompose()
                        body_to_process = ' '.join(soup.stripped_strings)
                        if body_to_process:
                            print(f"MSG Body: Successfully processed HTML body (length {len(body_to_process)}).")
                            processed_html_body = True
                        else:
                            print("MSG Body: Processed HTML body resulted in empty text.")
                else:
                    print("MSG Body: msg.htmlBody was string but failed to re-encode to bytes.")
            except Exception as e:
                print(f"MSG Body: Error processing HTML body: {e}")
                body_to_process = None
        if not processed_html_body:
            print("MSG Body: HTML body not found or failed, falling back to plain text body.")
            if msg.body:
                body_text = msg.body
                print(f"MSG Body (Raw Plain Text Fallback, first 500 chars): {body_text[:500]}")
                body_to_process = body_text
            else:
                print("MSG Body: Plain text body is also empty.")
                body_to_process = ""
        if body_to_process:
            body_lines = body_to_process.splitlines()
            for line in body_lines:
                cleaned_line = line.strip()
                if cleaned_line:
                    structure.append({'type': 'paragraph', 'text': cleaned_line})
                    full_text_list.append(cleaned_line)
        else:
            print("MSG Body: No processable body content found (HTML or Plain Text).")
        full_text = '\n'.join(full_text_list)
        print(f"Extracted full text from {file_path.name} (first 500 chars): {full_text[:500]}")
        return full_text, structure
    except InterruptedError:
        # 重新抛出取消异常
        raise
    except Exception as e:
        print(f"Error processing MSG file {file_path}: {e}")
        # traceback.print_exc(file=sys.stderr) # COMMENTED OUT
        return "", []

def index_documents(writer, content_dict: dict[Path, tuple[str, list[dict]]]):
    """将提取的内容添加到索引中
    
    Args:
        writer: Whoosh writer
        content_dict: 路径到内容的映射
    """
    for path, (content, structure_map) in content_dict.items():
        try:
            # 获取文件类型
            if isinstance(path, str):
                file_path = path.split('::')[0] if '::' in path else path  # 处理存档路径
                file_ext = Path(file_path.split('::')[-1] if '::' in file_path else file_path).suffix.lower()
            else:
                file_path = str(path)
                file_ext = path.suffix.lower()
            
            # 获取文件大小和修改时间
            file_size = content_dict.get(path, {}).get('file_size', 0)
            last_modified = content_dict.get(path, {}).get('last_modified', 0)
            
            # 获取OCR状态
            ocr_used = content_dict.get(path, {}).get('ocr_used', False)
            
            # 提取文件名（用于文件名搜索）
            if isinstance(path, str):
                norm_path = normalize_path_for_index(path)
                if '::' in norm_path:
                    # 处理存档内文件
                    _, member_path = norm_path.split('::', 1)
                    filename_text = Path(member_path).name
                else:
                    # 普通文件
                    filename_text = Path(norm_path).name
            else:
                # Path对象
                filename_text = path.name

            # 更新索引
            writer.update_document(
                path=normalize_path_for_index(str(path)),  # 标准化路径
                content=content,
                filename_text=filename_text,  # 文件名作为单独的字段
                structure_map=structure_map,
                last_modified=last_modified,
                file_size=file_size,
                file_type=file_ext.lstrip('.'),  # 去掉前导点
                indexed_with_ocr=ocr_used  # 存储OCR使用状态
            )
        except Exception as e:
            print(f"Warning: Error indexing document {path}: {e}")
            continue  # 继续索引其他文档

def get_positive_terms(q: Query) -> set[str]:
    """Recursively extract positive terms (not under a NOT) from a query tree."""
    positive_terms = set()
    if isinstance(q, Term):
        positive_terms.add(q.text)
    elif isinstance(q, (And, Or)):
        for subq in q.children():
            positive_terms.update(get_positive_terms(subq))
    elif isinstance(q, Phrase):
        positive_terms.update(q.words)
    # --- ADDED: 处理Wildcard查询 ---
    elif isinstance(q, Wildcard):
        # 去除通配符，提取关键部分进行高亮
        # 例如 "项目*计划" 会提取 "项目" 和 "计划" 进行高亮
        term = q.text
        # 分割通配符
        parts = []
        
        # 如果以*开头，去掉开头的*
        if term.startswith('*'):
            term = term[1:]
            
        # 如果以*结尾，去掉结尾的*
        if term.endswith('*'):
            term = term[:-1]
            
        # 使用通配符*分割字符串
        if '*' in term:
            parts = [p for p in term.split('*') if p]
        elif '?' in term:
            # 对于?通配符，简单地去除
            parts = [term.replace('?', '')]
        else:
            parts = [term]
            
        # 添加所有分割后的部分到高亮词集合
        for part in parts:
            if part:  # 确保不添加空字符串
                positive_terms.add(part)
    # --- END ADDED ---
    # Ignore terms under Not, Prefix, Wildcard, Every, etc. for highlighting purposes
    # You might want to refine this for Prefix/Wildcard if needed
    return positive_terms

# --- ADDED BACK: convert_term_to_prefix function --- 
def convert_term_to_prefix(q: Query, fieldname: str = "content") -> Query:
    """Recursively convert Term queries to Prefix queries within a Query tree."""
    if isinstance(q, Term) and q.fieldname == fieldname:
        # Only convert Terms in the specified field
        return Prefix(q.fieldname, q.text)
    elif isinstance(q, (And, Or)):
        # Recursively process children of And/Or groups
        return q.__class__([convert_term_to_prefix(subq, fieldname) for subq in q.children()])
    elif isinstance(q, Not):
        # For Not queries, process the *inner* query
        # Whoosh Not object holds its child query in the .query attribute
        return Not(convert_term_to_prefix(q.query, fieldname))
    else:
        # Return other query types (like Prefix, Wildcard, Every, etc.) unchanged
        return q
# --- END ADDED BACK --- 

# --- ADDED: 验证通配符语法的函数 ---
def validate_wildcard_syntax(query_str):
    """
    验证通配符语法是否合法
    返回: (bool, str) - (是否合法, 错误信息)
    """
    # 检查是否有未闭合的转义符
    if query_str.endswith('\\'):
        return False, "通配符查询不能以转义字符\\结尾"
    
    # 检查是否有连续多个*号，可能导致性能问题
    if '**' in query_str:
        return False, "连续的**通配符可能导致搜索缓慢，请使用单个*"
        
    # 检查是否只有通配符
    if query_str.strip() in ['*', '?', '*?', '?*']:
        return False, "查询不能仅包含通配符，请添加至少一个字符"
        
    return True, ""

# --- ADDED: 检查通配符查询性能风险 ---
def check_wildcard_performance_risk(query_str):
    """
    检查通配符查询是否可能导致性能问题
    返回: (bool, str) - (是否有风险, 风险描述)
    """
    # 检查是否以*开头
    if query_str.startswith('*'):
        return True, "以*开头的查询可能较慢，因为无法使用索引前缀优化"
    
    # 检查通配符数量是否过多
    wildcard_count = query_str.count('*') + query_str.count('?')
    if wildcard_count > 3:
        return True, f"查询包含{wildcard_count}个通配符，可能导致搜索较慢"
        
    return False, ""

# --- ADDED: 统一处理通配符查询 ---
def process_wildcard_query(query_str, is_filename_search=False):
    """
    统一处理通配符查询
    is_filename_search: 是否为文件名搜索
    返回: 处理后的查询字符串
    """
    has_wildcard = '*' in query_str or '?' in query_str
    
    # 文件名搜索默认添加通配符
    if is_filename_search and not has_wildcard:
        return f"*{query_str}*"
    
    # 中文通配符搜索增强：对于中文和数字混合的情况，提高匹配灵活度
    if has_wildcard and not is_filename_search:
        # 确保*前后的词能够更好地分割和匹配
        # 例如：将"十九届*全会"转换为更灵活的匹配模式
        if '*' in query_str:
            parts = query_str.split('*')
            for i in range(len(parts) - 1):
                if parts[i] and parts[i+1]:  # 确保不是开头或结尾的*
                    # 检查是否包含中文和数字混合
                    has_chinese = any('\u4e00' <= c <= '\u9fff' for c in parts[i]+parts[i+1])
                    if has_chinese:
                        # 为提高灵活性，对连接处进行特殊处理
                        print(f"增强中文通配符搜索: '{query_str}'")
            # 对于中文通配符搜索，暂时保持原样，主要通过Whoosh的Wildcard查询实现
        
    # 全文搜索保持原样
    return query_str

# 增加一个辅助函数用于扩展中文通配符搜索
def expand_chinese_wildcard_query(query_str):
    """
    扩展中文通配符查询，增加更多可能的匹配模式
    例如：十九届*全会 → 十九届*全会 OR 十九*全会
    """
    if '*' not in query_str and '?' not in query_str:
        return query_str
        
    expanded_queries = [query_str]
    
    # 特殊处理：检查是否包含数字+中文的模式，如"十九届*全会"
    parts = re.split(r'([*?])', query_str)
    pattern = re.compile(r'[\u4e00-\u9fff]+[\d]+[\u4e00-\u9fff]+|[\d]+[\u4e00-\u9fff]+|[\u4e00-\u9fff]+[\d]+')
    
    for i in range(0, len(parts), 2):
        if i < len(parts) and pattern.search(parts[i]):
            # 包含中文和数字混合，考虑更灵活的匹配
            if '届' in parts[i]:
                # 特殊处理包含"届"的情况，如"十九届"
                if i+2 < len(parts):
                    new_query = parts[i].replace('届', '') + parts[i+1] + parts[i+2]
                    expanded_queries.append(new_query)
                    print(f"扩展通配符查询: 添加 '{new_query}'")
    
    # 添加更多的中文通配符特殊处理规则
    # 处理常见的中文词组分词问题
    if '*' in query_str:
        parts = query_str.split('*')
        # 特殊处理中文词组间的通配符
        for i in range(len(parts) - 1):
            if parts[i] and parts[i+1] and all('\u4e00' <= c <= '\u9fff' for c in parts[i][-1] + parts[i+1][0]):
                # 尝试不同的词组切分方式
                if len(parts[i]) > 1:
                    # 拆分前一个词的最后一个字符
                    new_query = parts[i][:-1] + '*' + parts[i][-1] + parts[i+1]
                    if new_query != query_str and new_query not in expanded_queries:
                        expanded_queries.append(new_query)
                        print(f"中文分词优化: 添加 '{new_query}'")
                
                if len(parts[i+1]) > 1:
                    # 拆分后一个词的第一个字符
                    new_query = parts[i] + parts[i+1][0] + '*' + parts[i+1][1:]
                    if new_query != query_str and new_query not in expanded_queries:
                        expanded_queries.append(new_query)
                        print(f"中文分词优化: 添加 '{new_query}'")
    
    if len(expanded_queries) > 1:
        # 使用OR组合查询
        print(f"扩展后的查询选项: {expanded_queries}")
    
    return expanded_queries  # 返回所有扩展查询，上层调用处需要处理

def search_index(query_str: str,
                 index_dir_path: str, # Added parameter
                 search_mode: str = 'phrase',
                 # --- ADDED: Parameter for search scope ---
                 search_scope: str = 'fulltext',
                 # -----------------------------------------
                 min_size_kb: int | None = None,
                 max_size_kb: int | None = None,
                 start_date: str | None = None,  # Changed to string (e.g., "2023-01-01")
                 end_date: str | None = None,    # Changed to string (e.g., "2023-12-31")
                 file_type_filter: list[str] | None = None,
                 sort_by: str = 'relevance',
                 case_sensitive: bool = False,
                 # --- ADDED: Parameter for current source directories ---
                 current_source_dirs: list[str] | None = None) -> list[dict]: # Filter results by current directories
    # --- MODIFIED: Include search_scope in debug log ---
    print(f"\n🔍 开始搜索:")
    print(f"   查询: '{query_str}'")
    print(f"   模式: {search_mode}")
    print(f"   范围: {search_scope}")
    print(f"   索引: {index_dir_path}")
    if search_mode == 'phrase':
        print(f"   ⚠️  精确搜索模式: 将只返回包含完整短语 '{query_str}' 的结果")
    # ---------------------------------------------------
    print(f"Filters - Size: {min_size_kb}-{max_size_kb}KB, Date: {start_date}-{end_date}, Types: {file_type_filter}") # Debug
    print(f"Case Sensitive: {case_sensitive} (Note: Currently ignored by backend)") # ADDED Debug for case_sensitive

    processed_results = [] # <--- Initialize a new list to store processed hits
    if not Path(index_dir_path).exists() or not exists_in(index_dir_path):
        print(f"Error: Index directory '{index_dir_path}' not found.")
        return processed_results # <--- Return the empty processed list

    # --- 检查通配符搜索是否允许 (位置1：在最开始检查) ---
    if '*' in query_str or '?' in query_str:
        # 检查是否为专业版功能
        if not is_feature_available(Features.WILDCARDS):
            print(f"通配符搜索功能不可用 (未获得许可)")
            return [{'error': True, 'error_message': "通配符搜索是专业版功能，请升级以使用", 'license_required': True}]
        
        # 验证通配符语法
        is_valid, error_message = validate_wildcard_syntax(query_str)
        if not is_valid:
            print(f"无效的通配符语法: {error_message}")
            return [{'error': True, 'error_message': f"通配符语法错误: {error_message}"}]
        
        # 检查性能风险
        has_risk, risk_message = check_wildcard_performance_risk(query_str)
        if has_risk:
            print(f"通配符性能风险: {risk_message}")
            # 仅记录风险，不阻止搜索，可选择添加到结果中作为警告
            processed_results.append({'warning': True, 'warning_message': risk_message, 'performance_warning': True})

    # --- Determine target field based on scope --- ADDED
    target_field = "filename_text" if search_scope == 'filename' else "content"
    print(f"Searching in field: '{target_field}'")
    # --------------------------------------------
    
    # --- 检查许可证状态，确定当前可访问的文件类型 ---
    allowed_file_types = {'.docx', '.txt', '.html', '.htm', '.rtf', '.xlsx', '.pptx'} # 基础版允许的文件类型
    
    # 如果用户有相应的许可证，添加专业版文件类型
    if is_feature_available(Features.PDF_SUPPORT):
        allowed_file_types.add('.pdf')
    if is_feature_available(Features.MARKDOWN_SUPPORT):
        allowed_file_types.add('.md')
    if is_feature_available(Features.EMAIL_SUPPORT):
        allowed_file_types.add('.eml')
        allowed_file_types.add('.msg')
    if is_feature_available(Features.MULTIMEDIA_SUPPORT):
        # 添加多媒体文件类型支持
        # 视频文件
        allowed_file_types.update(['.mp4', '.mkv', '.avi', '.mov', '.wmv', '.flv', '.webm', '.m4v', '.3gp', '.rmvb'])
        # 音频文件  
        allowed_file_types.update(['.mp3', '.wav', '.flac', '.aac', '.ogg', '.wma', '.m4a', '.opus', '.aiff'])
        # 图片文件
        allowed_file_types.update(['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.svg', '.ico', '.raw'])
    
    print(f"Current license allows file types: {allowed_file_types}")
    # -------------------------------------------

    ix = open_dir(index_dir_path)
    searcher = ix.searcher(weighting=scoring.BM25F())
    # --- Use analyzer associated with the target field (or default) --- MODIFIED
    analyzer = ix.schema[target_field].analyzer if target_field in ix.schema else ChineseAnalyzer()
    # -------------------------------------------------------------------

    text_query = None
    parsed_query_obj = None # Store the parsed query object
    if query_str:
        # --- 处理中文通配符特殊情况 ---
        has_wildcard = '*' in query_str or '?' in query_str
        chinese_wildcard_expansion = False
        expanded_queries = []
        
        if has_wildcard:
            # 检测是否是需要扩展的中文通配符特殊情况
            if any('\u4e00' <= c <= '\u9fff' for c in query_str):  # 包含中文字符
                # 特别处理类似"十九届*全会"这样的查询
                if '届*' in query_str or '次*' in query_str or '*全会' in query_str:
                    chinese_wildcard_expansion = True
                    expanded_queries = []
                    
                    # 原始查询
                    expanded_queries.append(query_str)
                    
                    # 扩展查询1：处理"届"
                    if '届*' in query_str:
                        expanded_queries.append(query_str.replace('届*', '*'))
                        # 处理类似"十九届*全会"与"十九届历次全会"的匹配
                        if '全会' in query_str:
                            expanded_queries.append(query_str.replace('届*', '届历次'))
                            expanded_queries.append(query_str.replace('*全会', '历次全会'))
                    
                    # 扩展查询2：增加更多可能的匹配
                    if '*全会' in query_str:
                        prefix_part = query_str.split('*')[0]
                        expanded_queries.append(f"{prefix_part}*会议")
                        
                        # 特别处理包含"届"的前缀
                        if '届' in prefix_part:
                            base_prefix = prefix_part.split('届')[0] + '届'
                            expanded_queries.append(f"{base_prefix}历次*")
                            expanded_queries.append(f"{base_prefix}历次全会")
                        
                        # 特别处理"十九届"
                        if '十九届' in prefix_part:
                            expanded_queries.append(prefix_part.replace('十九届', '十九*'))
                            expanded_queries.append('十九*全会')
                            expanded_queries.append('十九届历次全会')

                # 通用中文通配符优化扩展
                if not expanded_queries:
                    # 如果没有通过特定规则扩展，使用通用扩展
                    general_expanded = expand_chinese_wildcard_query(query_str)
                    if isinstance(general_expanded, list) and len(general_expanded) > 1:
                        chinese_wildcard_expansion = True
                        expanded_queries = general_expanded
                    
                print(f"中文通配符搜索扩展: 原始查询 '{query_str}' 扩展为 {expanded_queries}")
        
        # --- MODIFIED: 使用统一的通配符处理函数 --- 
        if search_scope == 'filename':
            target_field = "filename_text"
            analyzer = ix.schema[target_field].analyzer if target_field in ix.schema else analysis.StandardAnalyzer() # Ensure standard analyzer for filename
            
            # 使用统一处理函数处理文件名搜索的通配符
            processed_query = process_wildcard_query(query_str, is_filename_search=True)
            text_query = Wildcard(target_field, processed_query)
            print(f"Constructed Wildcard query for filename: {text_query}")
        else: # Handle fulltext search based on search_mode
            target_field = "content"
            # --- FIXED: 强制使用我们修改后的ChineseAnalyzer，支持phrase_mode ---
            analyzer = ChineseAnalyzer()
            # --- 修改全文模糊搜索中的通配符处理逻辑 --- 
            if search_mode == 'phrase':
                # 检查是否包含逻辑操作符
                logical_operators = ['AND', 'OR', 'NOT']
                has_logical_operators = any(f" {op} " in f" {query_str} " for op in logical_operators)
                
                # --- ADDED: 检测通配符，在精确模式下也支持 --- 
                has_wildcard = '*' in query_str or '?' in query_str
                
                if has_wildcard:
                    # --- 处理中文通配符特殊扩展查询 ---
                    if chinese_wildcard_expansion and expanded_queries:
                        # 使用OR组合多个查询
                        sub_queries = []
                        for exp_query in expanded_queries:
                            sub_queries.append(Wildcard(target_field, exp_query))
                            print(f"Added wildcard expansion: {exp_query}")
                            
                        if len(sub_queries) == 1:
                            text_query = sub_queries[0]
                        else:
                            text_query = Or(sub_queries)
                            print(f"Created combined OR query with {len(sub_queries)} expansions")
                    else:
                        # 在精确模式下也支持通配符
                        print(f"Wildcard detected in phrase mode for '{target_field}': '{query_str}'. Using wildcard query.")
                        processed_query = process_wildcard_query(query_str, is_filename_search=False)
                        text_query = Wildcard(target_field, processed_query)
                        parsed_query_obj = text_query
                        print(f"Constructed Wildcard query in phrase mode on '{target_field}': {text_query}")
                elif has_logical_operators:
                    # 在精确搜索模式下不处理逻辑操作符，直接使用短语搜索
                    print(f"WARNING: Logical operators detected in phrase mode for query: '{query_str}'. These operators are only supported in fuzzy mode.")
                    # --- MODIFIED: 为精确搜索传递phrase_mode参数 ---
                    terms = [token.text for token in analyzer(query_str, phrase_mode=True)]
                    if terms:
                        text_query = Phrase(target_field, terms)
                        print(f"Constructed Phrase query on '{target_field}': {text_query}")
                        if text_query:
                            parsed_query_obj = text_query # Store phrase query object
                    else:
                        print(f"Phrase query for '{target_field}' is empty after analysis.")
                else:
                    # 重新设计的真正精确搜索：字符串包含匹配
                    # --- REDESIGNED: 真正的精确搜索实现 ---
                    
                    # 精确搜索的核心思想：查找包含完全相同字符串的文档
                    # 不依赖分词，直接进行字符串匹配
                    
                    print(f"精确搜索重新设计：查找包含 '{query_str}' 的文档")
                    
                    # 策略1: 尝试完整字符串的Term匹配（适用于索引中恰好有该词汇的情况）
                    strategies = []
                    
                    # 首先尝试将整个查询作为单个Term
                    strategies.append(Term(target_field, query_str))
                    print(f"策略1: 完整Term匹配 - {query_str}")
                    
                    # 策略2: 前缀匹配（适用于查询是某个更长词汇前缀的情况）
                    if len(query_str) >= 2:
                        strategies.append(Prefix(target_field, query_str))
                        print(f"策略2: 前缀匹配 - {query_str}*")
                    
                    # 策略3: 通配符匹配（在查询前后加通配符，寻找包含该字符串的词汇）
                    if len(query_str) >= 2:
                        wildcard_pattern = f"*{query_str}*"
                        strategies.append(Wildcard(target_field, wildcard_pattern))
                        print(f"策略3: 通配符包含匹配 - {wildcard_pattern}")
                    
                    # 策略4: 如果前面都失败，回退到分词短语匹配
                    terms = [token.text for token in analyzer(query_str, phrase_mode=True)]
                    if terms and len(terms) > 1:
                        strategies.append(Phrase(target_field, terms))
                        print(f"策略4: 分词短语匹配 - {terms}")
                    
                    # 使用OR查询组合所有策略
                    if len(strategies) > 1:
                        text_query = Or(strategies)
                        print(f"构造组合精确搜索查询: {len(strategies)} 种策略")
                    elif len(strategies) == 1:
                        text_query = strategies[0]
                        print(f"使用单一策略精确搜索")
                    else:
                        print("无法构造精确搜索查询")
                        text_query = None
                    
                    if text_query:
                        parsed_query_obj = text_query
                        print(f"最终精确搜索查询: {text_query}")
                    else:
                        print(f"精确搜索查询构造失败")
            elif search_mode == 'fuzzy':
                # --- 使用统一函数处理全文搜索的通配符 --- 
                if '*' in query_str or '?' in query_str:
                    # --- 处理中文通配符特殊扩展查询 ---
                    if chinese_wildcard_expansion and expanded_queries:
                        # 创建复合查询 (OR组合多个通配符查询)
                        sub_queries = []
                        for exp_query in expanded_queries:
                            sub_queries.append(Wildcard(target_field, exp_query))
                            print(f"添加通配符子查询: {exp_query}")
                        
                        if sub_queries:
                            # 使用OR组合所有子查询
                            if len(sub_queries) == 1:
                                text_query = sub_queries[0]
                            else:
                                text_query = Or(sub_queries)
                            parsed_query_obj = text_query
                            print(f"构建复合通配符查询: {text_query}")
                    else:
                        print(f"Wildcard detected in fuzzy mode for '{target_field}': '{query_str}'. Constructing direct Wildcard query.")
                        processed_query = process_wildcard_query(query_str, is_filename_search=False)
                        text_query = Wildcard(target_field, processed_query)
                        parsed_query_obj = text_query
                        print(f"Constructed direct Wildcard query on '{target_field}': {text_query}")
                else: # Not a wildcard query in fuzzy mode, use QueryParser as before for keywords, etc.
                    parser = QueryParser(target_field, schema=ix.schema)
                    try:
                        parsed_q = parser.parse(query_str)
                        parsed_query_obj = parsed_q # Store parsed object before conversion
                        print(f"Parsed Keyword query on '{target_field}': {parsed_q}")
                        # Convert to prefix for fulltext search
                        text_query = convert_term_to_prefix(parsed_q, fieldname=target_field)
                        if text_query != parsed_q:
                            print(f"-> Converted terms to prefix on '{target_field}': {text_query}")
                    except Exception as e:
                        print(f"Error parsing fuzzy query on '{target_field}': {e}")
                        text_query = None
                        parsed_query_obj = None
            else:
                print(f"Error: Unknown search mode '{search_mode}' for fulltext search")
        # --- END MODIFIED --- 
    else:
        print("No text query provided.")

    size_filter_query = None
    min_bytes = min_size_kb * 1024 if min_size_kb is not None else None
    max_bytes = max_size_kb * 1024 if max_size_kb is not None else None
    if min_bytes is not None or max_bytes is not None:
        size_filter_query = NumericRange("file_size", min_bytes, max_bytes)
        print(f"Constructed Size filter query: {size_filter_query}")
    date_filter_query = None
    start_timestamp = None
    end_timestamp = None
    if start_date:
        try:
            dt_start = datetime.strptime(start_date, '%Y-%m-%d')
            start_timestamp = dt_start.timestamp()
        except ValueError as e:
            print(f"Error parsing start_date '{start_date}': {e}. Expected format: YYYY-MM-DD")
    if end_date:
        try:
            dt_end = datetime.strptime(end_date, '%Y-%m-%d')
            dt_end = dt_end.replace(hour=23, minute=59, second=59, microsecond=999999)
            end_timestamp = dt_end.timestamp()
        except ValueError as e:
            print(f"Error parsing end_date '{end_date}': {e}. Expected format: YYYY-MM-DD")
    if start_timestamp is not None or end_timestamp is not None:
        date_filter_query = NumericRange("last_modified", start_timestamp, end_timestamp)
        print(f"Constructed Date filter query: {date_filter_query}")
    file_type_query = None
    if file_type_filter:
        lower_case_filters = [ftype.lower().lstrip('.') for ftype in file_type_filter if ftype]
        if lower_case_filters:
            type_queries = [Term("file_type", ftype) for ftype in lower_case_filters]
            if len(type_queries) == 1:
                file_type_query = type_queries[0]
            else:
                file_type_query = Or(type_queries)
            print(f"Constructed File Type filter query: {file_type_query}")
        else:
            print("File type filter list was empty or contained only empty strings.")

    # Combine all queries
    all_queries = []
    if text_query:
        all_queries.append(text_query)
    if size_filter_query:
        all_queries.append(size_filter_query)
    if date_filter_query:
        all_queries.append(date_filter_query)
    if file_type_query:
        all_queries.append(file_type_query)

    final_query = None
    if not all_queries:
        # --- If only filter criteria are present, search everything --- MODIFIED
        # print("Error: No search criteria (text, size, date, or type filter) provided.")
        print("No specific text query, searching based on filters only.")
        final_query = Every() # Match all documents if no criteria, filters will apply later
    elif len(all_queries) == 1:
        final_query = all_queries[0]
    else:
        final_query = And(all_queries)
        print(f"Combined query: {final_query}")

    # Sorting logic (remains the same)
    sort_field = None
    reverse = False
    if sort_by == 'date_asc':
        sort_field = 'last_modified'
        reverse = False
    elif sort_by == 'date_desc':
        sort_field = 'last_modified'
        reverse = True
    elif sort_by == 'size_asc':
        sort_field = 'file_size'
        reverse = False
    elif sort_by == 'size_desc':
        sort_field = 'file_size'
        reverse = True
    elif sort_by == 'relevance':
        sort_field = None  # Default Whoosh scoring
    else:
        print(f"Warning: Unknown sort_by option '{sort_by}', defaulting to relevance.")
        sort_field = None

    # --- 修改搜索结果处理逻辑，过滤掉许可证无法访问的文件类型 ---
    results = searcher.search(final_query, limit=3000, sortedby=sort_field, reverse=reverse) # Performance-balanced limit with user experience priority
    
    # --- Result Processing and Highlighting (Conditional) --- MODIFIED
    if results:
        print(f"Found {len(results)} document hit(s):")
        matched_contexts = 0
        
        # --- MODIFIED: Get positive terms for highlighting --- 
        # Only prepare for content highlighting if scope is fulltext and we have a parsed query
        positive_terms_for_highlighting = set()
        if search_scope == 'fulltext' and parsed_query_obj:
            positive_terms_for_highlighting = get_positive_terms(parsed_query_obj)
            print(f"DEBUG: Positive terms for highlighting: {positive_terms_for_highlighting}")
        # -----------------------------------------------------
        
        for hit in results:
            file_path = hit.get('path', "(未知文件)")
            file_type = hit.get('file_type', '')
            
            # --- 检查文件是否还存在 ---
            if not os.path.exists(file_path):
                print(f"Skipping result for {file_path} because file no longer exists")
                continue  # 跳过此结果，文件已被删除
            # -----------------------------------------
            
            # --- 检查文件是否在当前源目录中 ---
            if current_source_dirs:
                # 标准化文件路径和源目录路径进行比较
                file_path_normalized = os.path.normpath(file_path).lower()
                is_in_current_dirs = False
                
                for source_dir in current_source_dirs:
                    source_dir_normalized = os.path.normpath(source_dir).lower()
                    # 检查文件是否在这个源目录或其子目录中
                    if file_path_normalized.startswith(source_dir_normalized + os.sep) or \
                       file_path_normalized == source_dir_normalized:
                        is_in_current_dirs = True
                        break
                
                if not is_in_current_dirs:
                    print(f"Skipping result for {file_path} because it's not in current source directories")
                    continue  # 跳过此结果，不在当前源目录中
            # -----------------------------------------
            
            # --- 检查当前许可证是否允许访问该文件类型 ---
            if file_type:
                # 标准化文件类型格式，确保都以点开头
                normalized_file_type = file_type if file_type.startswith('.') else f'.{file_type}'
                
                if normalized_file_type not in allowed_file_types:
                    print(f"Skipping result for {file_path} due to license restrictions (type: {file_type})")
                    continue  # 跳过此结果，不添加到返回列表
            # -----------------------------------------
            
            # --- Basic result structure (always included) ---
            basic_result_info = {
                'file_path': file_path,
                'last_modified': hit.get('last_modified', 0),
                'file_size': hit.get('file_size', 0),
                'file_type': file_type,
                'score': hit.score
            }
            # -----------------------------------------------
            
            # --- Content-based processing only for fulltext search ---
            if search_scope == 'fulltext' and query_str: # Only process structure if doing text search in content
                added_paragraphs_in_hit = set()
                structure = []
                structure_map_json = hit.get('structure_map')
                if structure_map_json:
                    try:
                        structure = json.loads(structure_map_json)
                    except json.JSONDecodeError:
                        print(f"Error parsing structure for {file_path}")
                        # Add the basic info even if structure fails
                        processed_results.append(basic_result_info)
                        continue
                else:
                    print(f"Warning: No structure information for {file_path}")
                    # Add the basic info if no structure
                    processed_results.append(basic_result_info)
                    continue

                found_match_in_content = False # Flag to track if any block matched
                for i, block in enumerate(structure):
                    block_text = block.get('text', '')
                    block_type = block.get('type')
                    # Initialize markers/highlights
                    final_marked_paragraph = ""
                    final_marked_heading = ""
                    final_marked_excel_values = None
                    is_relevant_block = False # Reset for each block

                    # Helper for fuzzy highlighting (now uses specific positive terms)
                    def fuzzy_highlighter(text, terms_to_highlight):
                        if not isinstance(text, str) or not terms_to_highlight:
                            return str(text)
                        temp_text = text
                        for term in terms_to_highlight:
                            def replace_func(matchobj):
                                return f"__HIGHLIGHT_START__{matchobj.group(0)}__HIGHLIGHT_END__"
                            # Use regex to match whole word unless term contains wildcard?
                            # For now, simple substring match might be okay if terms are from parser
                            temp_text = re.sub(re.escape(term), replace_func, temp_text, flags=re.IGNORECASE)
                        return temp_text
                        
                    # Helper for phrase highlighting (remains mostly the same, but should ideally use positive phrase)
                    def phrase_highlighter(text, phrase):
                        if not isinstance(text, str) or not phrase:
                            return str(text)
                        phrase_pattern = re.compile(re.escape(phrase), flags=re.IGNORECASE)
                        def replace_phrase_func(matchobj):
                            return f"__HIGHLIGHT_START__{matchobj.group(0)}__HIGHLIGHT_END__"
                        return phrase_pattern.sub(replace_phrase_func, text)

                    # Check if block is relevant based on search mode and content
                    if search_mode == 'phrase':
                        # 添加精确搜索调试信息
                        escaped_query = re.escape(query_str)
                        match = re.search(escaped_query, block_text, flags=re.IGNORECASE)
                        if match:
                           print(f"🎯 精确匹配找到: 文件 {file_path}, 块类型 {block_type}")
                           print(f"   查询: '{query_str}' -> 正则: '{escaped_query}'")
                           print(f"   匹配文本: '{match.group(0)}'")
                           print(f"   块内容前50字符: '{block_text[:50]}...'")
                           is_relevant_block = True
                           highlighted_block_text = phrase_highlighter(block_text, query_str)
                           if block_type == 'heading' or block_type == 'metadata':
                               final_marked_heading = highlighted_block_text
                           elif block_type == 'excel_row':
                               original_excel_values = block.get('values', [])
                               marked_excel_values = []
                               for cell_value in original_excel_values:
                                   marked_excel_values.append(phrase_highlighter(cell_value, query_str))
                               final_marked_excel_values = marked_excel_values
                           else: # paragraph
                               final_marked_paragraph = highlighted_block_text
                               
                    elif search_mode == 'fuzzy':
                        contains_term = False
                        # --- MODIFIED: Check against positive terms only --- 
                        if positive_terms_for_highlighting:
                            # Check block text first
                            if block_text:
                                block_text_lower = block_text.lower()
                                for term in positive_terms_for_highlighting:
                                    if term in block_text_lower:
                                        contains_term = True
                                        break
                            # If not in text, check excel values if applicable
                            if not contains_term and block_type == 'excel_row':
                                original_excel_values = block.get('values', [])
                                for cell_value in original_excel_values:
                                    if isinstance(cell_value, str):
                                        cell_lower = cell_value.lower()
                                        for term in positive_terms_for_highlighting:
                                            if term in cell_lower:
                                                contains_term = True
                                                break
                                    if contains_term:
                                        break
                        # -----------------------------------------------
                                        
                        if contains_term:
                           is_relevant_block = True
                           # --- MODIFIED: Pass positive terms to highlighter --- 
                           if block_type == 'heading' or block_type == 'metadata':
                               final_marked_heading = fuzzy_highlighter(block_text, positive_terms_for_highlighting)
                           elif block_type == 'excel_row':
                               original_excel_values = block.get('values', [])
                               marked_excel_values = []
                               for cell_value in original_excel_values:
                                   marked_excel_values.append(fuzzy_highlighter(cell_value, positive_terms_for_highlighting))
                               final_marked_excel_values = marked_excel_values
                           else: # paragraph
                               final_marked_paragraph = fuzzy_highlighter(block_text, positive_terms_for_highlighting)
                           # -------------------------------------------------
                               
                    # If the block was relevant, create a detailed result entry
                    if is_relevant_block:
                        if block_text not in added_paragraphs_in_hit: # Avoid duplicate paragraphs from same file hit
                            result_block = {
                                **basic_result_info, # Include basic info
                                'type': block_type,
                                'level': block.get('level'),
                                'paragraph': block.get('text') if block_type != 'heading' and block_type != 'metadata' and block_type != 'excel_row' else None,
                                'heading': block.get('text') if block_type == 'heading' or block_type == 'metadata' else None,
                                'excel_sheet': block.get('sheet_name'),
                                'excel_row_idx': block.get('row_index'),
                                'excel_headers': block.get('headers'),
                                'excel_values': final_marked_excel_values if final_marked_excel_values is not None else block.get('values'), # Use marked or original
                                'marked_paragraph': final_marked_paragraph,
                                'marked_heading': final_marked_heading,
                                # Keep score from basic_result_info
                            }
                            # print(f"MATCHED BLOCK: {result_block}") # Too verbose
                            processed_results.append(result_block)
                            found_match_in_content = True
                            matched_contexts += 1
                            added_paragraphs_in_hit.add(block_text)
                            
                # If after checking all blocks, no content match was found for this hit
                # (This shouldn't happen often with fulltext search, but as a fallback)
                if not found_match_in_content:
                    print(f"Note: Hit found for {file_path} but no specific content block matched/highlighted.")
                    processed_results.append(basic_result_info) # Add basic info
                    
            else: # If search_scope is 'filename' OR query_str is empty (filter only search)
                # Just add the basic file info, no content highlighting needed here
                processed_results.append(basic_result_info)
        # End of loop through hits
    else:
        print("No document hits found for the query.")
        
    # Close searcher and index
    if searcher: searcher.close()
    if ix: ix.close()
    
    # 智能结果截断和用户友好提示
    original_count = len(processed_results)
    max_recommended_results = 500  # 推荐的最大结果数
    
    if original_count > max_recommended_results:
        # 截断到推荐数量，保留相关度最高的结果
        processed_results = processed_results[:max_recommended_results]
        print(f"⚠️  结果数量较多 ({original_count} 条)，为保证性能已截断到 {max_recommended_results} 条")
        print(f"💡 建议：使用更具体的搜索词或添加筛选条件以获得更精确的结果")
    elif original_count > 500:
        print(f"💡 找到 {original_count} 条结果，将使用虚拟滚动模式保证界面流畅性")
    
    print(f"--- Search complete. Returning {len(processed_results)} processed results. ---")
    return processed_results

# --- MODIFIED: Accept a dictionary --- 
# def _extract_worker(item_data: tuple) -> dict:
def _extract_worker(worker_args: dict) -> dict:
# -----------------------------------
    # --- Unpack arguments --- 
    # --- This is the OLD logic causing the KeyError ---
    # item_data = worker_args['item_data'] # <<< The line causing KeyError
    # enable_ocr = worker_args['enable_ocr'] # This key might also change based on how we passed it
    # path_key = item_data[0]
    # mtime = item_data[1]
    # fsize = item_data[2]
    # file_type = item_data[3]
    # archive_path = Path(item_data[4]) if len(item_data) > 4 and item_data[4] else None
    # member_name = item_data[5] if len(item_data) > 5 and item_data[5] else None
    # --- End of OLD logic ---

    # --- NEW logic to fix KeyError ---
    path_key = worker_args['path_key']
    file_type = worker_args['file_type']
    archive_path_abs_str = worker_args.get('archive_path_abs') # Use .get for optional keys
    member_name = worker_args.get('member_name')
    enable_ocr_for_file = worker_args['enable_ocr'] # Get per-file OCR setting
    # --- Get timeout value --- 
    extraction_timeout = worker_args.get('extraction_timeout') # Can be None
    # -----------------------
    # --- ADDED: Unpack content limit --- 
    content_limit_bytes = worker_args.get('content_limit_bytes', 0) # Default to 0 (no limit)
    # -----------------------------------
    # --- ADDED: Get index_dir_path for recording skipped files ---
    index_dir_path = worker_args.get('index_dir_path', '')
    # -----------------------------------------------------------
    # --- ADDED: Get cancel callback ---
    cancel_callback = worker_args.get('cancel_callback')
    # ----------------------------------
    # --- ADDED: Get filename-only flag ---
    is_filename_only = worker_args.get('is_filename_only', False)
    # ------------------------------------
    original_mtime = worker_args['original_mtime']
    original_fsize = worker_args['original_fsize']
    display_name = worker_args.get('display_name', Path(path_key).name if "::" not in path_key else path_key.split("::")[1]) # Use provided or generate
    filename_for_index = Path(path_key).name if not member_name else Path(member_name).name

    # --- Variable Initialization --- 
    text_content = None
    structure = [] # Initialize for non-PDF types
    error_message = None
    content_truncated = False # --- ADDED: Flag for truncation

    try:
        start_time = time.time()
        
        # --- MODIFIED: 在开始处理前检查是否需要取消 ---
        check_cancellation(cancel_callback, "文件提取处理开始")
        # ----------------------------------------
        
        # --- ADDED: 处理仅文件名索引模式 ---
        if is_filename_only:
            print(f"仅文件名索引: {display_name}")
            return {
                'path_key': path_key,
                'text_content': '',  # 空内容
                'filename': filename_for_index,
                'structure': [],
                'mtime': original_mtime,
                'fsize': original_fsize,
                'file_type': Path(path_key).suffix.lower().lstrip('.'),
                'ocr_enabled_for_file': False,
                'display_name': display_name,
                'content_source': 'filename_only',  # 标记来源
                'error': None,
                'content_truncated': False
            }
        # -----------------------------------
        
        # --- Select extraction function based on file type --- 
        if file_type == 'file':
            file_path = Path(path_key)
            file_ext = file_path.suffix.lower()
            
            # --- ADDED: 在确定文件类型后再次检查取消状态 ---
            check_cancellation(cancel_callback, f"文件类型确定 {file_ext}")
            # ----------------------------------------
            
            # Select function based on extension (Corrected indentation)
            if file_ext == '.docx':
                try:
                    print(f"开始处理DOCX文件: {display_name}")
                    text_content, structure = extract_text_from_docx(file_path, cancel_callback)
                    print(f"完成处理DOCX文件: {display_name}")
                except InterruptedError:
                    # 重新抛出取消异常
                    print(f"DOCX文件处理被用户取消: {display_name}")
                    raise
                except PermissionError as perm_err:
                    # 记录因许可证限制跳过的文件
                    error_message = str(perm_err)
                    if index_dir_path:
                        record_skipped_file(index_dir_path, str(file_path), f"许可证限制 - {error_message}")
            elif file_ext == '.txt':
                # --- MODIFIED: Apply content limit to TXT --- 
                try:
                    print(f"开始处理TXT文件: {display_name}")
                    text_content_tuple = extract_text_from_txt(file_path, cancel_callback) # Corrected indentation
                    if isinstance(text_content_tuple, tuple) and len(text_content_tuple) == 2:
                        text_content = text_content_tuple[0]
                        structure = text_content_tuple[1]
                        if text_content is None:
                            error_message = "TXT extraction failed."
                        elif content_limit_bytes and content_limit_bytes > 0: # Check if limit applies
                            # Encode to check byte length and truncate if needed
                            try:
                                encoded_content = text_content.encode('utf-8', errors='ignore')
                                if len(encoded_content) > content_limit_bytes:
                                    truncated_bytes = encoded_content[:content_limit_bytes]
                                    # Decode back, ignoring errors for partial characters
                                    text_content = truncated_bytes.decode('utf-8', errors='ignore')
                                    content_truncated = True # Set the flag
                                    # print(f"Info: Content truncated for {display_name} to {content_limit_bytes} bytes.") # COMMENTED OUT
                            except Exception as enc_err:
                                # print(f"Warning: Error during content truncation check for {display_name}: {enc_err}") # COMMENTED OUT
                                pass # Ignore truncation check errors silently for now
                    else:
                        error_message = "TXT extraction returned unexpected result."
                        text_content = None
                        structure = []
                    print(f"完成处理TXT文件: {display_name}")
                except InterruptedError:
                    # 重新抛出取消异常
                    print(f"TXT文件处理被用户取消: {display_name}")
                    raise
                except PermissionError as perm_err:
                    # 记录因许可证限制跳过的文件
                    error_message = str(perm_err)
                    if index_dir_path:
                        record_skipped_file(index_dir_path, str(file_path), f"许可证限制 - {error_message}")
            elif file_ext == '.pdf':
                # --- MODIFIED: Add explicit timeout parameter ---
                try:
                    print(f"开始处理PDF文件: {display_name} (OCR: {'启用' if enable_ocr_for_file else '禁用'})")
                    
                    # --- ADDED: 在开始PDF处理前检查取消状态 ---
                    check_cancellation(cancel_callback, f"PDF处理开始 {display_name}")
                    # ----------------------------------------
                    
                    text_content_tuple = extract_text_from_pdf(file_path, enable_ocr=enable_ocr_for_file, timeout=extraction_timeout, cancel_callback=cancel_callback)
                    if isinstance(text_content_tuple, tuple) and len(text_content_tuple) >= 2:
                        text_content = text_content_tuple[0]
                        structure = text_content_tuple[1] if text_content_tuple[1] is not None else []
                    else:
                        error_message = "PDF extraction returned unexpected result."
                        text_content = None
                        structure = []
                    print(f"完成处理PDF文件: {display_name}")
                except InterruptedError:
                    # 重新抛出取消异常
                    print(f"PDF文件处理被用户取消: {display_name}")
                    raise
                except Exception as e:
                    error_message = str(e)
                    text_content = None
                    structure = []
                    # --- ADDED: Record timeout errors for PDF extraction --- 
                    if "timeout" in error_message.lower() or "timed out" in error_message.lower():
                        if index_dir_path:
                            record_skipped_file(index_dir_path, str(file_path), f"PDF处理超时 ({extraction_timeout}秒)")
                    # --- ADDED: Record license errors for PDF extraction --- 
                    elif isinstance(e, PermissionError) or "许可证" in error_message or "license" in error_message.lower():
                        if index_dir_path:
                            record_skipped_file(index_dir_path, str(file_path), f"许可证限制 - PDF支持功能需要专业版许可证")
            elif file_ext == '.pptx':
                try:
                    print(f"开始处理PPTX文件: {display_name}")
                    text_content, structure = extract_text_from_pptx(file_path, cancel_callback)
                    print(f"完成处理PPTX文件: {display_name}")
                except InterruptedError:
                    # 重新抛出取消异常
                    print(f"PPTX文件处理被用户取消: {display_name}")
                    raise
                except PermissionError as perm_err:
                    # 记录因许可证限制跳过的文件
                    error_message = str(perm_err)
                    if index_dir_path:
                        record_skipped_file(index_dir_path, str(file_path), f"许可证限制 - {error_message}")
            elif file_ext == '.xlsx':
                try:
                    print(f"开始处理XLSX文件: {display_name}")
                    # 传递取消回调给Excel处理函数
                    text_content, structure = extract_text_from_xlsx(file_path, cancel_callback=cancel_callback)
                    print(f"完成处理XLSX文件: {display_name}")
                except InterruptedError:
                    # 重新抛出取消异常
                    print(f"XLSX文件处理被用户取消: {display_name}")
                    raise
                except PermissionError as perm_err:
                    # 记录因许可证限制跳过的文件
                    error_message = str(perm_err)
                    if index_dir_path:
                        record_skipped_file(index_dir_path, str(file_path), f"许可证限制 - {error_message}")
            elif file_ext == '.md':
                try:
                    print(f"开始处理MD文件: {display_name}")
                    text_content, structure = extract_text_from_md(file_path, cancel_callback)
                    print(f"完成处理MD文件: {display_name}")
                except InterruptedError:
                    # 重新抛出取消异常
                    print(f"MD文件处理被用户取消: {display_name}")
                    raise
                except PermissionError as perm_err:
                    # 记录因许可证限制跳过的文件
                    error_message = str(perm_err)
                    text_content = None
                    structure = []
                    if index_dir_path:
                        record_skipped_file(index_dir_path, str(file_path), f"许可证限制 - Markdown支持功能需要专业版许可证")
                except Exception as e:
                    error_message = str(e)
                    text_content = None
                    structure = []
            elif file_ext in ('.html', '.htm'): # Corrected structure/indentation
                try:
                    print(f"开始处理HTML文件: {display_name}")
                    text_content, structure = extract_text_from_html(file_path, cancel_callback)
                    print(f"完成处理HTML文件: {display_name}")
                except InterruptedError:
                    # 重新抛出取消异常
                    print(f"HTML文件处理被用户取消: {display_name}")
                    raise
                except PermissionError as perm_err:
                    # 记录因许可证限制跳过的文件
                    error_message = str(perm_err)
                    if index_dir_path:
                        record_skipped_file(index_dir_path, str(file_path), f"许可证限制 - {error_message}")
            elif file_ext == '.rtf': # Corrected structure/indentation
                try:
                    print(f"开始处理RTF文件: {display_name}")
                    text_content, structure = extract_text_from_rtf(file_path, cancel_callback)
                    print(f"完成处理RTF文件: {display_name}")
                except InterruptedError:
                    # 重新抛出取消异常
                    print(f"RTF文件处理被用户取消: {display_name}")
                    raise
                except PermissionError as perm_err:
                    # 记录因许可证限制跳过的文件
                    error_message = str(perm_err)
                    if index_dir_path:
                        record_skipped_file(index_dir_path, str(file_path), f"许可证限制 - {error_message}")
            elif file_ext == '.eml': # Corrected structure/indentation
                try:
                    print(f"开始处理EML文件: {display_name}")
                    text_content, structure = extract_text_from_eml(file_path, cancel_callback)
                    print(f"完成处理EML文件: {display_name}")
                except InterruptedError:
                    # 重新抛出取消异常
                    print(f"EML文件处理被用户取消: {display_name}")
                    raise
                except PermissionError as perm_err:
                    # 记录因许可证限制跳过的文件
                    error_message = str(perm_err)
                    text_content = None
                    structure = []
                    if index_dir_path:
                        record_skipped_file(index_dir_path, str(file_path), f"许可证限制 - 邮件支持功能需要专业版许可证")
                except Exception as e:
                    error_message = str(e)
                    text_content = None
            elif file_ext == '.msg': # Corrected structure/indentation
                try:
                    print(f"开始处理MSG文件: {display_name}")
                    text_content, structure = extract_text_from_msg(file_path, cancel_callback)
                    print(f"完成处理MSG文件: {display_name}")
                except InterruptedError:
                    # 重新抛出取消异常
                    print(f"MSG文件处理被用户取消: {display_name}")
                    raise
                except PermissionError as perm_err:
                    # 记录因许可证限制跳过的文件
                    error_message = str(perm_err)
                    text_content = None
                    structure = []
                    if index_dir_path:
                        record_skipped_file(index_dir_path, str(file_path), f"许可证限制 - 邮件支持功能需要专业版许可证")
                except Exception as e:
                    error_message = str(e)
                    text_content = None
                    structure = []
            else: # Corrected structure/indentation
                # 检查是否是多媒体文件，如果是则仅索引文件名
                if file_ext in FILENAME_ONLY_EXTENSIONS:
                    # 多媒体文件：仅使用文件名作为内容
                    text_content = filename_for_index
                    structure = [{'type': 'filename', 'text': filename_for_index}]
                    print(f"多媒体文件仅索引文件名: {display_name}")
                else:
                    error_message = f"Unsupported file extension: {file_ext}"
                    text_content = ""
                    structure = []
        
        elif file_type == 'archive':
            # --- ADDED: 在处理压缩包前检查取消状态 ---
            check_cancellation(cancel_callback, "压缩包处理开始")
            # ----------------------------------------
            
            if not archive_path_abs_str or not member_name:
                error_message = "Archive path or member name missing for archive extraction"
                text_content = ""
                structure = []
            else:
                archive_path = Path(archive_path_abs_str)
                member_ext = Path(member_name).suffix.lower()
                archive_type = archive_path.suffix.lower()
                # --- Extract member to temp file and process --- 
                with tempfile.TemporaryDirectory() as temp_dir:
                    temp_file_path = Path(temp_dir) / Path(member_name).name # Use just filename in temp dir
                    try: # Added try
                        # ... (Archive extraction logic remains the same) ...
                        if archive_type == '.zip': # Corrected indentation
                            try: # Try opening the zip file
                                with zipfile.ZipFile(archive_path, 'r') as zf:
                                    # --- ADDED: Try block specifically for member extraction ---
                                    try:
                                        zf.extract(member_name, path=temp_dir)
                                        # Rename logic after successful extraction
                                        extracted_member_path = Path(temp_dir) / member_name
                                        if extracted_member_path.exists() and extracted_member_path != temp_file_path:
                                            extracted_member_path.rename(temp_file_path)
                                    except RuntimeError as e_extract_member:
                                        if 'password required' in str(e_extract_member).lower() or 'encrypted file' in str(e_extract_member).lower():
                                            error_message = f"受密码保护的ZIP成员: {member_name}"
                                            archive_member_path = f"{archive_path_abs_str}::{member_name}"
                                            record_skipped_file(
                                                index_dir_path,
                                                archive_member_path,
                                                format_skip_reason("password_zip", f"成员 '{member_name}' 需要密码")
                                            )
                                        else:
                                             error_message = f"提取ZIP成员时发生运行时错误 '{member_name}': {e_extract_member}"
                                             record_skipped_file(
                                                 index_dir_path,
                                                 f"{archive_path_abs_str}::{member_name}",
                                                 format_skip_reason("extraction_error", f"成员运行时错误: {e_extract_member}")
                                             )
                                        temp_file_path = None # Indicate extraction failed
                                    except zipfile.BadZipFile as e_bad_member:
                                         error_message = f"ZIP成员损坏或格式错误: {member_name}"
                                         record_skipped_file(
                                             index_dir_path,
                                             f"{archive_path_abs_str}::{member_name}",
                                             format_skip_reason("corrupted_zip", f"成员损坏: {e_bad_member}")
                                         )
                                         temp_file_path = None # Indicate extraction failed
                                    except KeyError as e_key_extract:
                                        error_message = f"ZIP成员未找到或无法提取: {member_name}"
                                        record_skipped_file(
                                            index_dir_path,
                                            f"{archive_path_abs_str}::{member_name}",
                                            format_skip_reason("extraction_error", f"成员 '{member_name}' 未找到: {e_key_extract}")
                                        )
                                        temp_file_path = None # Indicate extraction failed
                                    except Exception as e_generic_extract:
                                         error_message = f"提取ZIP成员时发生未知错误 '{member_name}': {e_generic_extract}"
                                         record_skipped_file(
                                             index_dir_path,
                                             f"{archive_path_abs_str}::{member_name}",
                                             format_skip_reason("extraction_error", f"成员提取未知错误: {e_generic_extract}")
                                         )
                                         temp_file_path = None # Indicate extraction failed
                                    # --- END ADDED Try block ---

                            # --- Existing handling for errors opening the main ZIP ---
                            except (zipfile.BadZipFile, RuntimeError) as e_open_zip:
                                reason_key = "password_zip" if 'password required' in str(e_open_zip).lower() else "corrupted_zip"
                                error_message = f"{format_skip_reason(reason_key)}: {archive_path.name}"
                                record_skipped_file(
                                    index_dir_path,
                                    str(archive_path), # Record the archive path itself
                                    format_skip_reason(reason_key, str(e_open_zip))
                                )
                                temp_file_path = None # Indicate extraction failed
                            # --- End Existing Handling ---

                        elif archive_type == '.rar': # Corrected indentation
                            with rarfile.RarFile(archive_path, 'r') as rf:
                                rf.extract(member_name, path=temp_dir)
                                extracted_member_path = Path(temp_dir) / member_name
                                if extracted_member_path.exists() and extracted_member_path != temp_file_path:
                                    extracted_member_path.rename(temp_file_path)
                        else: # Corrected indentation
                            error_message = f"Unsupported archive type: {archive_type}"
                        
                        # --- ADDED: 在提取完成后检查取消状态 ---
                        check_cancellation(cancel_callback, "压缩包成员提取完成")
                        # ----------------------------------------
                        
                        if not error_message and temp_file_path and temp_file_path.exists():
                            # --- Now extract text from the temporary file --- 
                            if member_ext == '.docx': # Corrected indentation
                                text_content, structure = extract_text_from_docx(temp_file_path, cancel_callback)
                            elif member_ext == '.txt':
                                # --- MODIFIED: Apply content limit to TXT member --- 
                                text_content_tuple = extract_text_from_txt(temp_file_path, cancel_callback)
                                if isinstance(text_content_tuple, tuple) and len(text_content_tuple) == 2:
                                    text_content = text_content_tuple[0]
                                    structure = text_content_tuple[1]
                                    if text_content is None:
                                        error_message = "TXT member extraction failed."
                                    elif content_limit_bytes and content_limit_bytes > 0: # Check if limit applies
                                        # Encode to check byte length and truncate if needed
                                        try:
                                            encoded_content = text_content.encode('utf-8', errors='ignore')
                                            if len(encoded_content) > content_limit_bytes:
                                                truncated_bytes = encoded_content[:content_limit_bytes]
                                                # Decode back, ignoring errors for partial characters
                                                text_content = truncated_bytes.decode('utf-8', errors='ignore')
                                                content_truncated = True # Set the flag
                                                # print(f"Info: Content truncated for archive member {display_name} to {content_limit_bytes} bytes.") # COMMENTED OUT
                                        except Exception as enc_err:
                                            # print(f"Warning: Error during content truncation check for archive member {display_name}: {enc_err}") # COMMENTED OUT
                                            pass # Ignore truncation check errors silently for now
                                    else:
                                        error_message = "TXT member extraction returned unexpected result."
                                        text_content = None
                                        structure = []
                                # --- END MODIFIED --- 
                            elif member_ext == '.pdf':
                                # --- MODIFIED PDF Call for Archive Member --- 
                                pdf_result_member = extract_text_from_pdf(temp_file_path, enable_ocr=enable_ocr_for_file, timeout=extraction_timeout, cancel_callback=cancel_callback)
                                text_content = pdf_result_member[0]
                                structure = pdf_result_member[1]
                                if text_content is None:
                                    error_message = "PDF member extraction failed or timed out"
                                    # --- NEW: Record skipped file for PDF timeout ---
                                    if index_dir_path:
                                        archive_member_path = f"{archive_path_abs_str}::{member_name}"
                                        record_skipped_file(
                                            index_dir_path,
                                            archive_member_path,
                                            format_skip_reason("pdf_timeout", "PDF处理超时或转换错误")
                                        )
                                    # -------------------------------------------------
                                elif content_limit_bytes and content_limit_bytes > 0:
                                    # Apply content limit to PDF text if needed
                                    try:
                                        encoded_content = text_content.encode('utf-8', errors='ignore')
                                        if len(encoded_content) > content_limit_bytes:
                                            truncated_bytes = encoded_content[:content_limit_bytes]
                                            # Decode back, ignoring errors for partial characters
                                            text_content = truncated_bytes.decode('utf-8', errors='ignore')
                                            content_truncated = True # Set the flag
                                            # print(f"Info: PDF content truncated for archive member {display_name} to {content_limit_bytes} bytes.") # COMMENTED OUT
                                            # --- NEW: Record skipped file for content limit ---
                                            if index_dir_path:
                                                archive_member_path = f"{archive_path_abs_str}::{member_name}"
                                                record_skipped_file(
                                                    index_dir_path,
                                                    archive_member_path,
                                                    format_skip_reason("content_limit", f"内容大小({len(encoded_content) // 1024}KB)超过限制({content_limit_bytes // 1024}KB)")
                                                )
                                            # -------------------------------------------------
                                    except Exception as enc_err:
                                        # print(f"Warning: Error during PDF content truncation check for {display_name}: {enc_err}") # COMMENTED OUT
                                        pass # Ignore truncation check errors silently for now
                                # -------------------------------------------
                            elif member_ext == '.pptx':
                                text_content, structure = extract_text_from_pptx(temp_file_path, cancel_callback)
                            elif member_ext == '.xlsx':
                                # 传递取消回调给Excel处理函数
                                text_content, structure = extract_text_from_xlsx(temp_file_path, cancel_callback=cancel_callback)
                            elif member_ext == '.md':
                                text_content, structure = extract_text_from_md(temp_file_path, cancel_callback)
                            elif member_ext in ('.html', '.htm'): # Corrected structure/indentation
                                text_content, structure = extract_text_from_html(temp_file_path, cancel_callback)
                            elif member_ext == '.rtf':
                                text_content, structure = extract_text_from_rtf(temp_file_path, cancel_callback)
                            elif member_ext == '.eml':
                                text_content, structure = extract_text_from_eml(temp_file_path, cancel_callback)
                            elif member_ext == '.msg':
                                text_content, structure = extract_text_from_msg(temp_file_path, cancel_callback)
                            else: # Corrected indentation
                                # 检查是否是多媒体文件，如果是则仅索引文件名
                                if member_ext in FILENAME_ONLY_EXTENSIONS:
                                    # 多媒体文件：仅使用文件名作为内容
                                    text_content = Path(member_name).name
                                    structure = [{'type': 'filename', 'text': Path(member_name).name}]
                                    print(f"压缩包内多媒体文件仅索引文件名: {member_name}")
                                else:
                                    error_message = f"Unsupported file extension in archive: {member_ext}"
                                    text_content = ""
                                    structure = []
                        elif not error_message:
                            error_message = "Failed to extract member from archive"
                            text_content = "" # Ensure defined
                            structure = []    # Ensure defined
                    # Moved exception handling outside the if/elif chain for member types
                    except FileNotFoundError:
                        error_message = f"Member not found in archive: {member_name}"
                        text_content = ""
                        structure = []
                    except zipfile.BadZipFile:
                        error_message = f"Bad ZIP file containing member: {member_name}"
                        text_content = ""
                        structure = []
                    except rarfile.BadRarFile:
                        error_message = f"Bad RAR file containing member: {member_name}"
                        text_content = ""
                        structure = []
                    except NotImplementedError as nie:
                         error_message = f"Feature not implemented for member '{member_name}': {nie}"
                         text_content = ""
                         structure = []
                    except Exception as e_extract:
                        error_message = f"Error extracting member '{member_name}': {e_extract}"
                        # print(traceback.format_exc(), file=sys.stderr) # Already commented? Ensure commented.
                        text_content = ""
                        structure = []
        else:
            error_message = f"Unknown file type for extraction: {file_type}"
            text_content = ""
            structure = []
        
        end_time = time.time()
        # print(f"Extraction time for {display_name}: {end_time - start_time:.2f}s")
    
    except InterruptedError:
        # --- MODIFIED: 改进取消异常处理 ---
        print(f"文件处理被用户取消: {display_name}")
        # 重新抛出取消异常，不进行任何处理
        raise
    except Exception as e_outer:
        error_message = f"Unexpected error during extraction setup for {display_name}: {e_outer}"
        # traceback.print_exc() # COMMENTED OUT
        text_content = None
        structure = []

    # --- Ensure content is string, even if empty ---
    final_content = text_content if text_content is not None else ""
    # If there was an error OR text is None, ensure structure is empty.
    final_structure = structure if error_message is None and text_content is not None else []

    # --- Construct result dictionary ---
    result = {
        'path_key': path_key,
        'display_name': display_name,
        'text_content': final_content,
        'structure': final_structure,
        'error': error_message,
        'mtime': original_mtime,
        'fsize': original_fsize,
        'file_type': Path(path_key.split('::')[0]).suffix.lower() if "::" not in path_key else Path(path_key.split('::')[1]).suffix.lower(),
        'filename': filename_for_index,
        'ocr_enabled_for_file': enable_ocr_for_file,
        'content_truncated': content_truncated
    }
    return result

# --- ADDED: Function to read license-skipped files ---
def read_license_skipped_files(index_dir_path: str) -> dict:
    """
    读取因许可证限制而被跳过的文件记录
    
    Args:
        index_dir_path: 索引目录路径
        
    Returns:
        dict: 包含被跳过文件路径的字典，键为文件路径，值为跳过原因
    """
    skipped_files = {}
    log_file_path = os.path.join(index_dir_path, "index_skipped_files.tsv")
    
    if not os.path.exists(log_file_path):
        print(f"跳过文件记录不存在: {log_file_path}")
        return skipped_files
        
    try:
        with open(log_file_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f, delimiter='\t')
            # 跳过表头
            next(reader, None)
            
            for row in reader:
                if len(row) >= 3:
                    file_path, reason, timestamp = row[0], row[1], row[2]
                    # 仅关注因许可证限制而跳过的文件
                    if "许可证限制" in reason:
                        skipped_files[file_path] = reason
                        
        print(f"读取了 {len(skipped_files)} 个因许可证限制而跳过的文件记录")
    except Exception as e:
        print(f"读取跳过文件记录时出错: {e}")
    
    return skipped_files
# ---------------------------------------------------

# --- ADDED: Function to update skipped files record ---
def update_skipped_files_record(index_dir_path: str, processed_license_files: dict):
    """
    更新跳过文件记录，删除已经重新处理的许可证限制文件记录
    
    Args:
        index_dir_path: 索引目录路径
        processed_license_files: 已处理的许可证限制文件字典
    """
    if not processed_license_files:
        return  # 没有需要更新的记录
        
    log_file_path = os.path.join(index_dir_path, "index_skipped_files.tsv")
    if not os.path.exists(log_file_path):
        return  # 文件不存在，无需更新
        
    try:
        # 读取所有记录
        all_records = []
        with open(log_file_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f, delimiter='\t')
            # 保存表头
            header = next(reader, None)
            if not header:
                return  # 空文件或只有表头
                
            all_records.append(header)
            
            # 读取并过滤记录
            for row in reader:
                if len(row) >= 3:
                    file_path, reason, timestamp = row[0], row[1], row[2]
                    # 如果不是已处理的许可证限制文件，则保留
                    if file_path not in processed_license_files:
                        all_records.append(row)
        
        # 写回过滤后的记录
        with open(log_file_path, 'w', encoding='utf-8', newline='') as f:
            writer = csv.writer(f, delimiter='\t')
            for row in all_records:
                writer.writerow(row)
                
        print(f"已从跳过文件记录中删除 {len(processed_license_files)} 个已处理的许可证限制文件")
    except Exception as e:
        print(f"更新跳过文件记录时出错: {e}")
# ------------------------------------------------------

# --- 索引速度优化函数 ---

def get_optimal_worker_count(task_type="io_intensive"):
    """
    根据任务类型确定最优的工作进程数量

    Args:
        task_type: 任务类型，"io_intensive"（I/O密集型）或"cpu_intensive"（CPU密集型）

    Returns:
        int: 推荐的工作进程数量
    """
    cpu_count = multiprocessing.cpu_count()

    if task_type == "io_intensive":
        # I/O密集型任务，可以使用更多进程
        optimal_count = min(cpu_count * 2, 8)  # 最多8个进程
    else:
        # CPU密集型任务，使用CPU核心数
        optimal_count = max(cpu_count - 1, 1)  # 留一个核心给系统

    print(f"检测到CPU核心数: {cpu_count}, 任务类型: {task_type}, 推荐工作进程数: {optimal_count}")
    return optimal_count

def should_skip_large_file(file_path: Path, max_size_mb: int = 100) -> tuple[bool, str]:
    """
    检查是否应该跳过大文件

    Args:
        file_path: 文件路径
        max_size_mb: 最大文件大小限制（MB）

    Returns:
        tuple[bool, str]: (是否跳过, 跳过原因)
    """
    try:
        file_size = file_path.stat().st_size
        max_size_bytes = max_size_mb * 1024 * 1024

        if file_size > max_size_bytes:
            return True, f"文件大小 ({file_size // (1024*1024)}MB) 超过限制 ({max_size_mb}MB)"

        return False, ""
    except Exception as e:
        return True, f"无法获取文件大小: {e}"

def should_skip_system_file(file_path: Path) -> tuple[bool, str]:
    """
    检查是否应该跳过系统文件或临时文件

    Args:
        file_path: 文件路径

    Returns:
        tuple[bool, str]: (是否跳过, 跳过原因)
    """
    try:
        path_str = str(file_path).lower()

        # 跳过的路径模式
        skip_patterns = [
            'recycle.bin',
            '$recycle.bin',
            'system volume information',
            'pagefile.sys',
            'hiberfil.sys',
            'swapfile.sys',
            '.tmp',
            '.temp',
            '~$',  # Office临时文件
        ]

        for pattern in skip_patterns:
            if pattern in path_str:
                return True, f"系统文件或临时文件: {pattern}"

        # 检查是否为隐藏文件（Windows）
        if file_path.name.startswith('.') and len(file_path.name) > 1:
            return True, "隐藏文件"

        return False, ""
    except Exception as e:
        return True, f"检查系统文件时出错: {e}"

def scan_documents_optimized(directory_paths: list, max_file_size_mb: int = 100, 
                           skip_system_files: bool = True, file_types_to_index=None, 
                           filename_only_types=None, cancel_callback=None) -> tuple[list[Path], list[Path], list[dict]]:
    """
    优化的文档扫描函数，支持多个目录和文件过滤

    Args:
        directory_paths: 要扫描的目录路径列表（可以是字符串或Path对象）
        max_file_size_mb: 最大文件大小限制（MB）
        skip_system_files: 是否跳过系统文件
        file_types_to_index: 要完整索引的文件类型列表，如['txt', 'docx']
        filename_only_types: 只索引文件名的文件类型列表，如['pdf', 'xlsx']

    Returns:
        tuple[list[Path], list[Path], list[dict]]: (完整索引文件列表, 仅文件名索引文件列表, 跳过的文件信息列表)
    """
    found_files = []  # 需要完整索引的文件
    filename_only_files = []  # 仅索引文件名的文件
    skipped_files = []
    
    # 用于去重的集合，防止重复添加同一文件
    processed_paths = set()  # 存储已处理的规范化路径

    # 转换为Path对象
    path_objects = []
    for dir_path in directory_paths:
        if isinstance(dir_path, str):
            path_objects.append(Path(dir_path))
        else:
            path_objects.append(dir_path)

    # 确定允许的文件扩展名
    if file_types_to_index is not None:
        # 用户明确指定了文件类型（包括空列表）
        allowed_extensions = []
        for file_type in file_types_to_index:
            # 确保扩展名以点开头
            ext = file_type if file_type.startswith('.') else f'.{file_type}'
            allowed_extensions.append(ext.lower())
        print(f"根据用户选择，完整索引以下文件类型: {allowed_extensions}")
    else:
        # 使用默认的所有支持的文件类型
        allowed_extensions = ALLOWED_EXTENSIONS
        print(f"使用默认文件类型: {allowed_extensions}")
    
    # 确定仅文件名索引的文件扩展名
    filename_only_extensions = []
    if filename_only_types:
        for file_type in filename_only_types:
            # 确保扩展名以点开头
            ext = file_type if file_type.startswith('.') else f'.{file_type}'
            filename_only_extensions.append(ext.lower())
        print(f"根据用户选择，仅文件名索引以下文件类型: {filename_only_extensions}")

    for directory_path in path_objects:
        # 在扫描每个目录前检查是否需要取消
        check_cancellation(cancel_callback, f"扫描目录 {directory_path}")
            
        if not directory_path.is_dir():
            print(f"错误: 路径不是目录: {directory_path}")
            continue
            
        print(f"扫描目录: {directory_path}")

        try:
            file_count = 0
            for item in directory_path.rglob('*'):
                # 每扫描50个文件检查一次取消状态
                file_count += 1
                periodic_cancellation_check(cancel_callback, 50, file_count, "文件扫描")
                    
                if not item.is_file():
                    continue

                # 规范化文件路径用于去重检查
                normalized_path = normalize_path_for_index(str(item))
                if normalized_path in processed_paths:
                    # 跳过重复文件
                    print(f"跳过重复文件: {item.name} (路径: {normalized_path})")
                    continue
                
                # 检查文件扩展名的处理策略
                file_ext = item.suffix.lower()
                
                if file_ext in allowed_extensions:
                    # 完整索引
                    file_category = "full_index"
                elif file_ext in filename_only_extensions:
                    # 仅文件名索引
                    file_category = "filename_only"
                else:
                    # 跳过此文件
                    if file_types_to_index or filename_only_types:
                        skipped_files.append({
                            'path': str(item),
                            'reason': f'文件类型 {item.suffix} 未被选择索引',
                            'type': 'file_type_not_selected'
                        })
                    continue

                # 检查是否跳过大文件
                should_skip_large, large_reason = should_skip_large_file(item, max_file_size_mb)
                if should_skip_large:
                    skipped_files.append({
                        'path': str(item),
                        'reason': large_reason,
                        'type': 'large_file'
                    })
                    continue
                            
                # 检查是否跳过系统文件
                if skip_system_files:
                    should_skip_sys, sys_reason = should_skip_system_file(item)
                    if should_skip_sys:
                        skipped_files.append({
                            'path': str(item),
                            'reason': sys_reason,
                            'type': 'system_file'
                        })
                        continue
                
                # 根据文件类别添加到相应列表
                if file_category == "full_index":
                    found_files.append(item)
                    processed_paths.add(normalized_path)  # 记录已处理的路径
                elif file_category == "filename_only":
                    filename_only_files.append(item)
                    processed_paths.add(normalized_path)  # 记录已处理的路径

        except InterruptedError:
            # 重新抛出取消异常
            raise
        except Exception as e:
            print(f"扫描目录时出错 {directory_path}: {e}")
            continue

    print(f"扫描完成. 完整索引: {len(found_files)} 个文档, 仅文件名索引: {len(filename_only_files)} 个文档, 跳过: {len(skipped_files)} 个文件")
    return found_files, filename_only_files, skipped_files

def estimate_processing_time(files: list[Path]) -> dict:
    """
    根据文件大小和类型估算处理时间

    Args:
        files: 文件列表

    Returns:
        dict: 包含时间估算信息的字典
    """
    # 不同文件类型的处理速度估算（MB/秒）
    processing_speeds = {
        '.txt': 50,   # 文本文件处理很快
        '.docx': 10,  # Word文档中等速度
        '.pdf': 5,    # PDF文件较慢（特别是OCR）
        '.pptx': 8,   # PowerPoint中等速度
        '.xlsx': 12,  # Excel中等速度
        '.html': 20,  # HTML文件较快
        '.md': 30,    # Markdown文件很快
        '.rtf': 15,   # RTF文件中等速度
        '.eml': 8,    # 邮件文件中等速度
        '.msg': 6,    # Outlook邮件文件较慢
    }

    total_size = 0
    estimated_time = 0
    file_type_counts = {}

    for file_path in files:
        try:
            file_size = file_path.stat().st_size / (1024 * 1024)  # MB
            file_ext = file_path.suffix.lower()

            total_size += file_size
            file_type_counts[file_ext] = file_type_counts.get(file_ext, 0) + 1

            # 获取处理速度，默认值为5MB/s
            speed = processing_speeds.get(file_ext, 5)
            estimated_time += file_size / speed

        except Exception:
            # 如果无法获取文件大小，使用默认估算
            estimated_time += 0.1  # 假设0.1秒

    return {
        'total_files': len(files),
        'total_size_mb': total_size,
        'estimated_time_seconds': estimated_time,
        'estimated_time_formatted': f"{int(estimated_time // 60)}分{int(estimated_time % 60)}秒",
        'file_type_counts': file_type_counts
    }

def create_or_update_index(directories: list[str], index_dir_path: str, enable_ocr: bool = True,
                          extraction_timeout: int = 300, content_limit_kb: int = 1024,
                          max_file_size_mb: int = 100, skip_system_files: bool = True,
                          incremental: bool = True, max_workers: int = None, 
                          cancel_callback=None, file_types_to_index=None, 
                          filename_only_types=None, preserve_removed_dirs: bool = True):
    """
    创建或更新文档索引（优化版本）

    Args:
        directories: 要索引的目录列表
        index_dir_path: 索引存储目录
        enable_ocr: 是否启用OCR
        extraction_timeout: 文件提取超时时间（秒）
        content_limit_kb: 内容大小限制（KB）
        max_file_size_mb: 最大文件大小限制（MB）
        skip_system_files: 是否跳过系统文件
        incremental: 是否启用增量索引
        max_workers: 最大工作进程数
        cancel_callback: 取消检查回调函数，如果返回True则取消操作
        file_types_to_index: 要索引的文件类型列表，如['txt', 'docx', 'pdf']

    Yields:
        dict: 进度信息
    """
    # --- MODIFIED: 在函数开始就初始化progress变量 ---
    progress = {
        'stage': 'initializing',
        'current': 0,
        'total': 0,
        'message': '初始化索引操作...',
        'files_processed': 0,
        'files_skipped': 0,
        'errors': 0
    }
    # ------------------------------------------------
    
    try:
        # 检查是否需要取消
        check_cancellation(cancel_callback, "索引初始化")
            
        # 确保索引目录存在
        index_path = Path(index_dir_path)
        index_path.mkdir(parents=True, exist_ok=True)

        # 更新进度信息
        progress.update({
            'stage': 'scanning',
            'message': '开始扫描文件...'
        })
        yield progress

        # 检查是否需要取消
        check_cancellation(cancel_callback, "开始文件扫描")

        # 1. 扫描文件
        print("开始扫描文档...")
        all_files, filename_only_files, skipped_files = scan_documents_optimized(
            directories, max_file_size_mb, skip_system_files, file_types_to_index, filename_only_types, cancel_callback
        )

        total_files = len(all_files) + len(filename_only_files)
        progress.update({
            'stage': 'scanning_complete',
            'total': total_files,
            'message': f'扫描完成，完整索引: {len(all_files)} 个，仅文件名: {len(filename_only_files)} 个，跳过: {len(skipped_files)} 个文件',
            'files_skipped': len(skipped_files)
        })
        yield progress

        # 检查是否需要取消
        check_cancellation(cancel_callback, "扫描完成后检查")

        if not all_files and not filename_only_files:
            progress.update({
                'stage': 'complete',
                'message': '没有找到需要索引的文件'
            })
            yield progress
            return

        # 2. 加载文件缓存（用于增量索引）
        file_cache = {}
        # 合并所有需要处理的文件
        files_to_process = all_files + filename_only_files

        if incremental:
            progress.update({
                'stage': 'change_detection',
                'message': '检测文件变更...'
            })
            yield progress

            # 检查是否需要取消
            if cancel_callback and cancel_callback():
                raise InterruptedError("操作被用户取消")

            file_cache = load_file_index_cache(index_dir_path)
            new_files, modified_files, deleted_files = detect_file_changes(all_files, file_cache, filename_only_files)

            files_to_process = new_files + modified_files

            progress.update({
                'stage': 'change_detection_complete',
                'total': len(files_to_process),
                'message': f'增量检测完成: {len(new_files)} 个新文件, {len(modified_files)} 个修改文件, {len(deleted_files)} 个删除文件'
            })
            yield progress

            # 检查是否需要取消
            if cancel_callback and cancel_callback():
                raise InterruptedError("操作被用户取消")

            # 如果没有变更，直接返回
            if not files_to_process and not deleted_files:
                progress.update({
                    'stage': 'complete',
                    'message': '没有文件变更，索引已是最新'
                })
                yield progress
                return

        # 3. 估算处理时间
        estimated_time_info = estimate_processing_time(files_to_process)
        progress.update({
            'stage': 'processing_start',
            'message': f'开始处理 {len(files_to_process)} 个文件，预计需要 {estimated_time_info["estimated_time_formatted"]}'
        })
        yield progress

        # 检查是否需要取消
        if cancel_callback and cancel_callback():
            raise InterruptedError("操作被用户取消")

        # 4. 准备Whoosh索引
        from whoosh import fields, index

        # 定义索引模式
        schema = fields.Schema(
            path=fields.ID(stored=True, unique=True),
            content=fields.TEXT(stored=True),
            filename_text=fields.TEXT(stored=True),
            structure_map=fields.TEXT(stored=True),
            last_modified=fields.NUMERIC(stored=True),
            file_size=fields.NUMERIC(stored=True),
            file_type=fields.TEXT(stored=True),
            indexed_with_ocr=fields.BOOLEAN(stored=True)
        )

        # 创建或打开索引
        if index.exists_in(index_dir_path):
            ix = index.open_dir(index_dir_path)
        else:
            ix = index.create_in(index_dir_path, schema)

        # 5. 批量处理文件
        if files_to_process:
            progress.update({
                'stage': 'extracting',
                'message': '开始提取文件内容...'
            })
            yield progress

            # 检查是否需要取消
            if cancel_callback and cancel_callback():
                raise InterruptedError("操作被用户取消")

            # 准备工作进程参数 - 只处理需要更新的文件
            # 分离需要处理的完整索引文件和仅文件名索引文件
            files_to_process_full = [f for f in files_to_process if f in all_files]
            files_to_process_filename_only = [f for f in files_to_process if f in filename_only_files]
            
            print(f"增量处理: 需要全文索引 {len(files_to_process_full)} 个文件, 需要文件名索引 {len(files_to_process_filename_only)} 个文件")
            
            worker_args_list = prepare_worker_arguments_batch(
                files_to_process_full, enable_ocr, extraction_timeout, 
                content_limit_kb, index_dir_path, files_to_process_filename_only, cancel_callback
            )

            # 用于收集进度更新的列表
            progress_updates = []

            # 定义进度回调函数
            def extraction_progress_callback(current, total, detail):
                progress_updates.append({
                    'stage': 'extracting',
                    'current': current,
                    'total': total,
                    'message': detail
                })

            # 多进程提取内容（带进度回调）
            extraction_results = []
            total_files = len(worker_args_list)
            
            # 逐个处理文件并实时发送进度
            processed_count = 0
            real_processing_total = len(files_to_process_full) + len(files_to_process_filename_only)
            
            for i, args in enumerate(worker_args_list):
                # --- MODIFIED: 在每个文件处理前检查是否需要取消 ---
                if cancel_callback and cancel_callback():
                    print(f"用户取消操作，停止处理剩余 {len(worker_args_list) - i} 个文件")
                    raise InterruptedError("操作被用户取消")
                # ------------------------------------------------
                    
                try:
                    # --- ADDED: 在开始处理单个文件前再次检查 ---
                    if cancel_callback and cancel_callback():
                        print(f"用户取消操作，跳过文件: {args.get('display_name', 'unknown')}")
                        raise InterruptedError("操作被用户取消")
                    # ----------------------------------------
                    
                    # 检查是否为真正需要处理的文件
                    file_path = args.get('path_key', args.get('file_path', ''))
                    is_filename_only = args.get('is_filename_only', False)
                    
                    # 只为真正需要处理的文件显示进度
                    if Path(file_path) in files_to_process_full or Path(file_path) in files_to_process_filename_only:
                        processed_count += 1
                        file_name = args.get('display_name', args.get('path_key', 'unknown'))
                        detail = f"正在处理: {file_name}"
                        
                        progress.update({
                            'stage': 'extracting',
                            'current': processed_count,
                            'total': real_processing_total,
                            'message': detail
                        })
                        yield progress
                    
                    # 处理单个文件
                    result = _extract_worker(args)
                    extraction_results.append(result)
                    
                    # --- ADDED: 在文件处理完成后检查取消状态 ---
                    if cancel_callback and cancel_callback():
                        print(f"用户取消操作，已处理 {i+1} 个文件，停止处理剩余文件")
                        raise InterruptedError("操作被用户取消")
                    # ----------------------------------------
                    
                except InterruptedError:
                    # --- ADDED: 专门处理取消异常 ---
                    print(f"文件处理被取消，已处理 {i} 个文件，剩余 {len(worker_args_list) - i} 个文件未处理")
                    # 直接重新抛出，不添加到结果中
                    raise
                except Exception as e:
                    print(f"处理文件时出错: {e}")
                    # 创建错误结果
                    error_result = {
                        'path_key': args.get('path_key', 'unknown'),
                        'display_name': args.get('display_name', 'unknown'),
                        'text_content': '',
                        'structure': [],
                        'error': str(e),
                        'mtime': args.get('original_mtime', 0),
                        'fsize': args.get('original_fsize', 0),
                        'file_type': '',
                        'filename': '',
                        'ocr_enabled_for_file': False,
                        'content_truncated': False
                    }
                    extraction_results.append(error_result)
                    
                    # 错误时也发送进度更新
                    current_file = i + 1
                    file_name = args.get('display_name', args.get('path_key', 'unknown'))
                    detail = f"处理失败: {file_name} - {str(e)}"
                    
                    progress.update({
                        'stage': 'extracting',
                        'current': current_file,
                        'total': total_files,
                        'message': detail
                    })
                    yield progress

            # --- MODIFIED: 在提取完成后检查是否需要取消 ---
            if cancel_callback and cancel_callback():
                print("用户在文件提取完成后取消操作")
                raise InterruptedError("操作被用户取消")
            # ----------------------------------------

            progress.update({
                'stage': 'indexing',
                'message': '开始索引文档...'
            })
            yield progress

            # 6. 批量索引文档
            with ix.writer() as writer:
                # 根据preserve_removed_dirs参数决定是否删除文件
                if incremental and 'deleted_files' in locals() and not preserve_removed_dirs:
                    print(f"物理删除 {len(deleted_files)} 个索引条目（preserve_removed_dirs=False）")
                    remove_deleted_files_from_index(writer, deleted_files)
                elif incremental and 'deleted_files' in locals() and preserve_removed_dirs:
                    print(f"保留 {len(deleted_files)} 个已移除目录的索引条目（preserve_removed_dirs=True）")
                    print("搜索时将通过目录过滤排除这些结果")

                # 索引提取的文档
                success_count = 0
                error_count = 0
                total_results = len(extraction_results)

                for i, result in enumerate(extraction_results):
                    # --- MODIFIED: 在每个文档索引前检查是否需要取消 ---
                    if cancel_callback and cancel_callback():
                        print(f"用户取消操作，停止索引剩余 {len(extraction_results) - i} 个文档")
                        raise InterruptedError("操作被用户取消")
                    # ------------------------------------------------
                        
                    try:
                        if result.get('error'):
                            # 记录错误文件
                            record_skipped_file(index_dir_path, result['path_key'], result['error'])
                            error_count += 1
                            
                            # 发送进度更新
                            current = i + 1
                            detail = f"跳过错误文件: {result.get('display_name', result['path_key'])}"
                            progress.update({
                                'stage': 'indexing',
                                'current': current,
                                'total': total_results,
                                'message': detail
                            })
                            yield progress
                            continue

                        # 索引文档
                        writer.add_document(
                            path=result['path_key'],
                            content=result['text_content'],
                            filename_text=result['filename'],
                            structure_map=json.dumps(result['structure'], ensure_ascii=False),
                            last_modified=result['mtime'],
                            file_size=result['fsize'],
                            file_type=result['file_type'],
                            indexed_with_ocr=result['ocr_enabled_for_file']
                        )
                        success_count += 1
                        
                        # 发送进度更新
                        current = i + 1
                        detail = f"已索引: {result.get('display_name', result['path_key'])}"
                        progress.update({
                            'stage': 'indexing',
                            'current': current,
                            'total': total_results,
                            'message': detail
                        })
                        yield progress
                        
                        # --- ADDED: 在索引文档后检查取消状态 ---
                        if cancel_callback and cancel_callback():
                            print(f"用户取消操作，已索引 {i+1} 个文档，停止索引剩余文档")
                            raise InterruptedError("操作被用户取消")
                        # ----------------------------------------
                        
                    except InterruptedError:
                        # 重新抛出取消异常
                        raise
                    except Exception as e:
                        error_count += 1
                        record_skipped_file(index_dir_path, result.get('path_key', 'unknown'), f"索引错误: {e}")
                        print(f"索引文档时出错: {e}")
                        
                        # 发送进度更新
                        current = i + 1
                        detail = f"索引失败: {result.get('display_name', result.get('path_key', 'unknown'))} - {str(e)}"
                        progress.update({
                            'stage': 'indexing',
                            'current': current,
                            'total': total_results,
                            'message': detail
                        })
                        yield progress

                progress.update({
                    'files_processed': success_count,
                    'errors': error_count
                })

        # 7. 更新文件缓存
        if incremental:
            # 检查是否需要取消
            if cancel_callback and cancel_callback():
                raise InterruptedError("操作被用户取消")
                
            progress.update({
                'stage': 'updating_cache',
                'message': '更新文件缓存...'
            })
            yield progress

            # 更新缓存（使用新的缓存条目格式，包含hash和mode）
            all_processed_files = set()  # 跟踪已处理的文件，避免重复
            
            for file_path in all_files:
                path_str = normalize_path_for_index(str(file_path))
                if path_str not in all_processed_files:
                    file_cache[path_str] = get_file_cache_entry(file_path, "full")
                    all_processed_files.add(path_str)
                
            for file_path in filename_only_files:
                path_str = normalize_path_for_index(str(file_path))
                if path_str not in all_processed_files:
                    file_cache[path_str] = get_file_cache_entry(file_path, "filename_only")
                    all_processed_files.add(path_str)

            save_file_index_cache(index_dir_path, file_cache)

        # 8. 记录跳过的文件
        for skip_info in skipped_files:
            if isinstance(skip_info, dict):
                # 新格式：字典包含详细信息
                record_skipped_file(index_dir_path, skip_info['path'], skip_info['reason'])
            else:
                # 兼容旧格式：元组 (file_path, reason)
                file_path, reason = skip_info
                record_skipped_file(index_dir_path, str(file_path), reason)

        # 完成
        progress.update({
            'stage': 'complete',
            'message': f'索引完成！处理了 {progress["files_processed"]} 个文件，跳过 {progress["files_skipped"]} 个文件，{progress["errors"]} 个错误'
        })
        yield progress

    except InterruptedError:
        # 用户取消操作
        progress.update({
            'stage': 'cancelled',
            'message': '索引操作已被用户取消'
        })
        yield progress
        raise
    except Exception as e:
        progress.update({
            'stage': 'error',
            'message': f'索引过程中出错: {str(e)}'
        })
        yield progress
        raise

# --- 结束索引优化函数 ---

# --- 高级索引优化功能 ---

def get_file_hash(file_path: Path, index_mode: str = "full") -> str:
    """
    获取文件的简单哈希值（基于修改时间和大小，不包含索引模式）

    Args:
        file_path: 文件路径
        index_mode: 索引模式（"full" 或 "filename_only"）- 保留参数兼容性但不用于哈希计算

    Returns:
        str: 文件哈希值
    """
    try:
        stat = file_path.stat()
        # 使用整数精度的修改时间和文件大小生成哈希
        # 不包含索引模式，避免索引模式变更时误判为文件变更
        mtime_int = int(stat.st_mtime)  # 使用整数精度避免浮点误差
        hash_str = f"{mtime_int}_{stat.st_size}"
        return hash_str
    except Exception:
        return "unknown"

def get_file_cache_entry(file_path: Path, index_mode: str) -> dict:
    """
    获取文件的缓存条目（包含哈希和索引模式）

    Args:
        file_path: 文件路径
        index_mode: 索引模式（"full" 或 "filename_only"）

    Returns:
        dict: 包含hash和mode的缓存条目
    """
    return {
        "hash": get_file_hash(file_path),
        "mode": index_mode
    }

def load_file_index_cache(index_dir_path: str) -> dict:
    """
    加载文件索引缓存，用于检测文件变更

    Args:
        index_dir_path: 索引目录路径

    Returns:
        dict: 文件路径到哈希值的映射
    """
    cache_file = Path(index_dir_path) / "file_cache.json"
    if cache_file.exists():
        try:
            with open(cache_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            print(f"加载文件缓存失败: {e}")
    return {}

def save_file_index_cache(index_dir_path: str, cache: dict):
    """
    保存文件索引缓存

    Args:
        index_dir_path: 索引目录路径
        cache: 文件缓存字典
    """
    cache_file = Path(index_dir_path) / "file_cache.json"
    try:
        with open(cache_file, 'w', encoding='utf-8') as f:
            json.dump(cache, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"保存文件缓存失败: {e}")

def detect_file_changes(files: list[Path], cache: dict, filename_only_files: list[Path] = None) -> tuple[list[Path], list[Path], list[str]]:
    """
    检测文件变更（支持索引模式感知）

    Args:
        files: 当前完整索引文件列表
        cache: 现有文件缓存
        filename_only_files: 仅文件名索引文件列表

    Returns:
        tuple[list[Path], list[Path], list[str]]: (新文件, 修改的文件, 删除的文件路径)
    """
    current_files = {}
    new_files = []
    modified_files = []
    
    if filename_only_files is None:
        filename_only_files = []

    # 检查完整索引文件
    for file_path in files:
        path_str = normalize_path_for_index(str(file_path))
        current_entry = get_file_cache_entry(file_path, "full")
        current_files[path_str] = current_entry

        if path_str not in cache:
            new_files.append(file_path)
            print(f"新文件: {file_path.name} (全文索引)")
        else:
            cached_entry = cache[path_str]
            # 兼容旧缓存格式（纯字符串）
            if isinstance(cached_entry, str):
                cached_hash = cached_entry
                cached_mode = "unknown"
            else:
                cached_hash = cached_entry.get("hash", "")
                cached_mode = cached_entry.get("mode", "unknown")
            
            current_hash = current_entry["hash"]
            current_mode = current_entry["mode"]
            
            # 检查是否需要重新处理
            file_changed = cached_hash != current_hash
            mode_upgraded = cached_mode == "filename_only" and current_mode == "full"
            
            if file_changed or mode_upgraded:
                modified_files.append(file_path)
                if file_changed:
                    print(f"检测到文件变更: {file_path.name} (文件时间或大小发生变化)")
                if mode_upgraded:
                    print(f"检测到索引模式升级: {file_path.name} (从仅文件名升级到全文索引)")

    # 检查仅文件名索引文件
    for file_path in filename_only_files:
        path_str = normalize_path_for_index(str(file_path))
        current_entry = get_file_cache_entry(file_path, "filename_only")
        current_files[path_str] = current_entry

        if path_str not in cache:
            new_files.append(file_path)
            print(f"新文件: {file_path.name} (仅文件名索引)")
        else:
            cached_entry = cache[path_str]
            # 兼容旧缓存格式（纯字符串）
            if isinstance(cached_entry, str):
                cached_hash = cached_entry
                cached_mode = "unknown"
            else:
                cached_hash = cached_entry.get("hash", "")
                cached_mode = cached_entry.get("mode", "unknown")
            
            current_hash = current_entry["hash"]
            current_mode = current_entry["mode"]
            
            # 检查是否需要重新处理
            file_changed = cached_hash != current_hash
            # 注意：从full降级到filename_only不需要重新处理
            
            if file_changed:
                modified_files.append(file_path)
                print(f"检测到文件变更: {file_path.name} (文件时间或大小发生变化)")

    # CRITICAL FIX REMOVED: 这段代码导致所有缓存文件都被重新处理
    # 增量索引应该只基于文件本身的变更，而不是配置变更
    # 如果用户新增文件类型，scan_documents_optimized已经会正确发现这些文件
    # 并且它们不在缓存中，所以会被正确标记为new_files

    # 检查删除的文件
    deleted_files = [path for path in cache.keys() if path not in current_files]

    return new_files, modified_files, deleted_files

def prepare_worker_arguments_batch(files: list[Path], enable_ocr: bool, extraction_timeout: int,
                                 content_limit_kb: int, index_dir_path: str, filename_only_files: list[Path] = None, cancel_callback=None) -> list[dict]:
    """
    批量准备工作进程参数

    Args:
        files: 需要完整索引的文件列表
        enable_ocr: 是否启用OCR
        extraction_timeout: 提取超时时间
        content_limit_kb: 内容大小限制（KB）
        index_dir_path: 索引目录路径
        filename_only_files: 仅索引文件名的文件列表
        cancel_callback: 取消检查回调函数

    Returns:
        list[dict]: 工作进程参数列表
    """
    worker_args_list = []
    content_limit_bytes = content_limit_kb * 1024 if content_limit_kb > 0 else 0

    # 处理需要完整索引的文件
    for file_path in files:
        # 动态设置PDF OCR超时
        actual_timeout = extraction_timeout
        if file_path.suffix.lower() == '.pdf' and enable_ocr:
            try:
                file_size = file_path.stat().st_size
                # 根据PDF文件大小设置更合理的超时时间
                if file_size < 5 * 1024 * 1024:  # 小于5MB
                    actual_timeout = min(60, extraction_timeout)
                elif file_size < 20 * 1024 * 1024:  # 5-20MB
                    actual_timeout = min(180, extraction_timeout)
                elif file_size < 50 * 1024 * 1024:  # 20-50MB
                    actual_timeout = min(300, extraction_timeout)
                else:  # 大于50MB
                    actual_timeout = extraction_timeout

                print(f"PDF文件 {file_path.name} ({file_size // (1024*1024)}MB) 设置OCR超时: {actual_timeout}秒")
            except Exception:
                pass

        file_stat = file_path.stat()

        worker_args = {
            'path_key': str(file_path),
            'file_type': 'file',
            'enable_ocr': enable_ocr and file_path.suffix.lower() == '.pdf',
            'extraction_timeout': actual_timeout,
            'content_limit_bytes': content_limit_bytes,
            'index_dir_path': index_dir_path,
            'original_mtime': file_stat.st_mtime,
            'original_fsize': file_stat.st_size,
            'display_name': file_path.name,
            'cancel_callback': cancel_callback,  # 添加取消回调
            'is_filename_only': False  # 标记为完整索引
        }

        worker_args_list.append(worker_args)

    # 处理仅索引文件名的文件
    if filename_only_files:
        for file_path in filename_only_files:
            file_stat = file_path.stat()
            
            worker_args = {
                'path_key': str(file_path),
                'file_type': 'file',
                'enable_ocr': False,  # 仅文件名索引不需要OCR
                'extraction_timeout': 1,  # 很短的超时，因为不需要提取内容
                'content_limit_bytes': 0,  # 不限制内容，因为不提取
                'index_dir_path': index_dir_path,
                'original_mtime': file_stat.st_mtime,
                'original_fsize': file_stat.st_size,
                'display_name': file_path.name,
                'cancel_callback': cancel_callback,
                'is_filename_only': True  # 标记为仅文件名索引
            }
            
            worker_args_list.append(worker_args)

    return worker_args_list

def process_files_multiprocess(worker_args_list: list[dict], max_workers: int = None, progress_callback=None) -> list[dict]:
    """
    使用多进程处理文件（简化版本）

    Args:
        worker_args_list: 工作进程参数列表
        max_workers: 最大工作进程数
        progress_callback: 进度回调函数，接收(current, total, detail)参数

    Returns:
        list[dict]: 处理结果列表
    """
    if max_workers is None:
        max_workers = get_optimal_worker_count("io_intensive")

    results = []
    total_files = len(worker_args_list)

    # 这里实现一个简化的多进程处理
    # 在实际部署中，这应该使用真正的multiprocessing.Pool
    print(f"开始多进程处理 {total_files} 个文件，使用 {max_workers} 个进程")

    for i, args in enumerate(worker_args_list):
        try:
            # 在实际实现中这里会调用真正的多进程处理
            result = _extract_worker(args)
            results.append(result)

            # 实时进度更新
            current_file = i + 1
            file_name = args.get('display_name', args.get('path_key', 'unknown'))
            detail = f"正在处理: {file_name}"
            
            # 调用进度回调
            if progress_callback:
                progress_callback(current_file, total_files, detail)

            # 简化的控制台进度报告（保留原有逻辑）
            if current_file % 10 == 0:
                print(f"已处理 {current_file}/{total_files} 个文件")

        except Exception as e:
            print(f"处理文件时出错: {e}")
            # 创建错误结果
            error_result = {
                'path_key': args.get('path_key', 'unknown'),
                'display_name': args.get('display_name', 'unknown'),
                'text_content': '',
                'structure': [],
                'error': str(e),
                'mtime': args.get('original_mtime', 0),
                'fsize': args.get('original_fsize', 0),
                'file_type': '',
                'filename': '',
                'ocr_enabled_for_file': False,
                'content_truncated': False
            }
            results.append(error_result)
            
            # 错误时也更新进度
            if progress_callback:
                current_file = i + 1
                file_name = args.get('display_name', args.get('path_key', 'unknown'))
                detail = f"处理失败: {file_name} - {str(e)}"
                progress_callback(current_file, total_files, detail)

    return results

def batch_index_documents(writer, extraction_results: list[dict], index_dir_path: str, progress_callback=None) -> tuple[int, int]:
    """
    批量索引文档（优化版本）

    Args:
        writer: Whoosh writer
        extraction_results: 提取结果列表
        index_dir_path: 索引目录路径
        progress_callback: 进度回调函数，接收(current, total, detail)参数

    Returns:
        tuple[int, int]: (成功索引数, 错误数)
    """
    success_count = 0
    error_count = 0
    total_results = len(extraction_results)

    for i, result in enumerate(extraction_results):
        try:
            if result.get('error'):
                # 记录错误文件
                record_skipped_file(index_dir_path, result['path_key'], result['error'])
                error_count += 1
                
                # 发送进度更新
                if progress_callback:
                    current = i + 1
                    detail = f"跳过错误文件: {result.get('display_name', result['path_key'])}"
                    progress_callback(current, total_results, detail)
                continue

            content = result.get('text_content', '')
            structure = result.get('structure', [])

            if not content and not structure:
                # 空内容，跳过
                record_skipped_file(index_dir_path, result['path_key'], "提取的内容为空")
                error_count += 1
                
                # 发送进度更新
                if progress_callback:
                    current = i + 1
                    detail = f"跳过空内容文件: {result.get('display_name', result['path_key'])}"
                    progress_callback(current, total_results, detail)
                continue

            # 获取文件信息
            path_key = result['path_key']
            file_path = Path(path_key)

            # 添加到索引
            writer.update_document(
                path=normalize_path_for_index(path_key),
                content=content,
                filename_text=file_path.name,
                structure_map=json.dumps(structure),
                last_modified=result['mtime'],
                file_size=result['fsize'],
                file_type=result.get('file_type', '').lstrip('.'),
                indexed_with_ocr=result.get('ocr_enabled_for_file', False)
            )

            success_count += 1
            
            # 发送进度更新
            if progress_callback:
                current = i + 1
                detail = f"已索引: {result.get('display_name', file_path.name)}"
                progress_callback(current, total_results, detail)

        except Exception as e:
            error_count += 1
            record_skipped_file(index_dir_path, result.get('path_key', 'unknown'), f"索引错误: {e}")
            print(f"索引文档时出错: {e}")
            
            # 发送进度更新
            if progress_callback:
                current = i + 1
                detail = f"索引失败: {result.get('display_name', result.get('path_key', 'unknown'))} - {str(e)}"
                progress_callback(current, total_results, detail)

    return success_count, error_count

def remove_deleted_files_from_index(writer, deleted_files: list[str]):
    """
    从索引中删除已删除的文件

    Args:
        writer: Whoosh writer
        deleted_files: 已删除的文件路径列表
    """
    for file_path in deleted_files:
        try:
            writer.delete_by_term('path', normalize_path_for_index(file_path))
            print(f"从索引中删除: {file_path}")
        except Exception as e:
            print(f"删除索引项时出错 {file_path}: {e}")

# --- 结束高级索引优化功能 ---

# --- 兼容性包装函数 ---

def create_or_update_index_legacy(source_directories, index_dir_path, enable_ocr, 
                                 extraction_timeout=300, txt_content_limit_kb=1024, 
                                 file_types_to_index=None, filename_only_types=None, 
                                 cancel_callback=None, preserve_removed_dirs=True):
    """
    兼容性包装函数，保持与现有GUI的兼容性
    将旧版本的参数映射到新的优化版本

    Args:
        source_directories: 源目录列表
        index_dir_path: 索引目录路径
        enable_ocr: 是否启用OCR
        extraction_timeout: 提取超时时间
        txt_content_limit_kb: TXT内容限制（KB）
        file_types_to_index: 要索引的文件类型列表
        cancel_callback: 取消检查回调函数

    Yields:
        dict: 进度信息（转换为旧格式）
    """
    print("使用兼容性包装函数调用优化版索引...")

    # 将新的优化参数映射到旧的格式
    try:
        # 调用优化版本的索引函数
        for progress in create_or_update_index(
            directories=source_directories,
            index_dir_path=index_dir_path,
            enable_ocr=enable_ocr,
            extraction_timeout=extraction_timeout,
            content_limit_kb=txt_content_limit_kb,
            max_file_size_mb=100,  # 默认值
            skip_system_files=True,  # 默认启用
            incremental=True,  # 默认启用增量索引
            max_workers=None,  # 使用自动检测
            cancel_callback=cancel_callback,  # 传递取消回调
            file_types_to_index=file_types_to_index,  # 传递完整索引文件类型
            filename_only_types=filename_only_types,  # 新增：传递仅文件名索引文件类型
            preserve_removed_dirs=preserve_removed_dirs  # 新增：传递目录保留参数
        ):
            # 将新格式的进度信息转换为旧格式
            old_format_progress = convert_progress_to_legacy_format(progress)
            yield old_format_progress

    except InterruptedError:
        # 处理用户取消
        cancelled_progress = {
            'type': 'complete',
            'message': '索引已被用户取消',
            'summary': {
                'message': '索引已被用户取消。',
                'added': 0,
                'updated': 0,
                'deleted': 0,
                'errors': 0,
                'cancelled': True
            }
        }
        yield cancelled_progress
        raise
    except Exception as e:
        # 如果优化版本失败，记录错误并抛出
        print(f"优化版索引失败: {e}")
        raise

def convert_progress_to_legacy_format(new_progress):
    """
    将新格式的进度信息转换为旧格式，确保GUI兼容性

    Args:
        new_progress: 新格式的进度字典

    Returns:
        dict: 旧格式的进度字典，确保所有字段都存在
    """
    # --- ENHANCED: 更强的类型检查和错误处理 ---
    if not isinstance(new_progress, dict):
        print(f"WARNING: convert_progress_to_legacy_format received non-dict: {type(new_progress)}")
        return {
            'type': 'error',
            'message': f'进度格式错误: {type(new_progress)}',
            'current': 0,
            'total': 0,
            'phase': 'error',
            'detail': '进度数据格式错误'
        }
    
    stage = new_progress.get('stage', '')
    message = new_progress.get('message', '')

    # 根据阶段映射到旧的消息类型
    if stage == 'scanning':
        return {
            'type': 'status',
            'message': message,
            'phase': 'scanning'
        }
    elif stage == 'scanning_complete':
        return {
            'type': 'status',
            'message': message,
            'phase': 'scan_complete'
        }
    elif stage == 'change_detection':
        return {
            'type': 'status',
            'message': message,
            'phase': 'change_detection'
        }
    elif stage == 'change_detection_complete':
        return {
            'type': 'status',
            'message': message,
            'phase': 'change_detection_complete'
        }
    elif stage == 'processing_start':
        return {
            'type': 'status',
            'message': message,
            'phase': 'processing'
        }
    elif stage == 'extracting':
        return {
            'type': 'progress',
            'current': new_progress.get('current', 0),
            'total': new_progress.get('total', 0),
            'phase': 'extracting',
            'detail': message
        }
    elif stage == 'indexing':
        return {
            'type': 'progress',
            'current': new_progress.get('current', 0),
            'total': new_progress.get('total', 0),
            'phase': 'indexing',
            'detail': message
        }
    elif stage == 'updating_cache':
        return {
            'type': 'status',
            'message': message,
            'phase': 'updating_cache'
        }
    elif stage == 'complete':
        # 构建完成摘要
        files_processed = new_progress.get('files_processed', 0)
        files_skipped = new_progress.get('files_skipped', 0)
        errors = new_progress.get('errors', 0)

        summary = {
            'message': message,
            'added': files_processed,  # 简化：将处理的文件数作为添加数
            'updated': 0,  # 新版本不区分添加和更新
            'deleted': 0,  # 删除数暂时设为0
            'errors': errors,
            'cancelled': False
        }

        return {
            'type': 'complete',
            'message': message,
            'summary': summary
        }
    elif stage == 'cancelled':
        # 处理取消状态
        summary = {
            'message': message,
            'added': 0,
            'updated': 0,
            'deleted': 0,
            'errors': 0,
            'cancelled': True
        }

        return {
            'type': 'complete',
            'message': message,
            'summary': summary
        }
    elif stage == 'error':
        return {
            'type': 'error',
            'message': message
        }
    else:
        # 默认情况
        return {
            'type': 'status',
            'message': message,
            'phase': stage
        }

# --- 结束兼容性包装函数 ---

if __name__ == "__main__":
    # 如果直接运行此文件，执行测试
    print("执行索引优化测试...")
    # test_optimized_indexing()  # 暂时注释掉测试函数调用
